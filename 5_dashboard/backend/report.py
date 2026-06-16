"""
1페이지 PPTX 보고서 생성.
레이아웃: 좌(Modeling 현황) / 우(불량 Unit 분석 현황)
디자인: 흰 배경 라이트 테마
이미지박스와 텍스트박스는 완전히 분리된 독립 영역.
"""
import io
import os
from datetime import datetime, timedelta
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

from graphs import ensure_all_graphs

# ── 색상 (라이트 테마) ────────────────────────────────────────
C_HEADER   = RGBColor(0x1E, 0x3A, 0x5F)  # 헤더 진한 파랑
C_TAB      = RGBColor(0x2A, 0x4A, 0x6F)  # 섹션 탭
C_SECTION  = RGBColor(0xF0, 0xF4, 0xF8)  # 섹션 배경 (연한 회색)
C_CARD     = RGBColor(0xFF, 0xFF, 0xFF)  # 카드 흰색
C_IMGBOX   = RGBColor(0xE8, 0xF0, 0xFE)  # 이미지 플레이스홀더
C_BTN      = RGBColor(0x3B, 0x82, 0xF6)  # Position 버튼 파랑
C_WAFER    = RGBColor(0x60, 0xA5, 0xFA)  # 웨이퍼맵
C_RED      = RGBColor(0xEF, 0x44, 0x44)  # 강조 빨강
C_GREEN    = RGBColor(0x16, 0xA3, 0x4A)  # 긍정 초록
C_WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
C_TEXT     = RGBColor(0x1E, 0x29, 0x3B)  # 본문 텍스트
C_SUB      = RGBColor(0x64, 0x74, 0x8B)  # 보조 텍스트
C_BORDER   = RGBColor(0xCB, 0xD5, 0xE1)  # 테두리
C_FOOTER   = RGBColor(0xF1, 0xF5, 0xF9)  # 푸터 배경
C_SUBTEXT  = RGBColor(0xBF, 0xDB, 0xFF)  # 헤더 내 보조

W = Inches(13.33)
H = Inches(7.5)


def _new_prs():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs


def _rect(slide, x, y, w, h, fill, line_color=None, line_pt=0.5):
    shape = slide.shapes.add_shape(1, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line_pt)
    else:
        shape.line.fill.background()
    return shape


def _oval(slide, x, y, w, h, fill):
    shape = slide.shapes.add_shape(9, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    return shape


def _rounded(slide, x, y, w, h, fill):
    shape = slide.shapes.add_shape(5, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    return shape


def _text(slide, txt, x, y, w, h, size=10, bold=False,
          color=None, align=PP_ALIGN.LEFT, italic=False):
    txb = slide.shapes.add_textbox(x, y, w, h)
    tf  = txb.text_frame
    tf.word_wrap = True
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = txt
    run.font.size   = Pt(size)
    run.font.bold   = bold
    run.font.italic = italic
    run.font.color.rgb = color or C_TEXT
    return txb


def _multiline(slide, lines, x, y, w, h, size=10, color=None,
               bold=False, line_spacing=1.1):
    txb = slide.shapes.add_textbox(x, y, w, h)
    tf  = txb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p   = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run()
        run.text = line
        run.font.size  = Pt(size)
        run.font.color.rgb = color or C_TEXT
        run.font.bold  = bold
    return txb


def _card(slide, x, y, w, h, lines=None, size=10, bold=False,
          title=None, title_size=10):
    """흰 카드 박스 (테두리 있음)."""
    _rect(slide, x, y, w, h, C_CARD, line_color=C_BORDER, line_pt=0.75)
    cy = y + Inches(0.08)
    if title:
        _text(slide, title, x + Inches(0.1), cy, w - Inches(0.2), Inches(0.28),
              size=title_size, bold=True, color=C_TAB)
        cy += Inches(0.28)
    if lines:
        _multiline(slide, lines,
                   x + Inches(0.1), cy,
                   w - Inches(0.2), h - (cy - y) - Inches(0.05),
                   size=size, color=C_TEXT, bold=bold)


def _imgbox(slide, x, y, w, h, label, img_path=None):
    """
    이미지 플레이스홀더 또는 실제 이미지 삽입.
    img_path가 있으면 실제 이미지, 없으면 빈 박스.
    """
    if img_path and os.path.exists(img_path):
        slide.shapes.add_picture(img_path, x, y, w, h)
        # 테두리용 투명 rect 덮기
        border = slide.shapes.add_shape(1, x, y, w, h)
        border.fill.background()
        border.line.color.rgb = C_BORDER
        border.line.width = Pt(0.75)
    else:
        _rect(slide, x, y, w, h, C_IMGBOX, line_color=C_BORDER, line_pt=0.75)
        _text(slide, label,
              x, y + h // 2 - Inches(0.18), w, Inches(0.36),
              size=9, color=C_SUB, align=PP_ALIGN.CENTER, italic=True)


# ── matplotlib 차트 → PNG bytes ──────────────────────────────
def _chart_trend_png(lot_labels, lot_production, lot_pred_yield, w_px=580, h_px=120, defect_ppm=None, ww_labels_in=None) -> bytes:
    """L2 불량 트렌드: 생산량 막대 + 예측불량ppm 꺾은선."""
    import matplotlib
    matplotlib.use("Agg")
    import matplotlib.pyplot as plt
    import matplotlib.font_manager as fm
    import numpy as np

    _try_set_font()

    n = len(lot_labels)
    if n == 0:
        lot_labels    = [f"WW{37-i}" for i in range(7, 0, -1)]
        lot_production = [1800, 2000, 1600, 2100, 1900, 1950, 1750]
        defect_ppm     = [105053, 433934, 233063, 336283, 210989, 290582, 23747]
        n = len(lot_labels)

    # defect_ppm 우선, 없으면 pred_yield(수율%)에서 역산
    if defect_ppm and len(defect_ppm) == n:
        ppm = [round(float(v)) if v is not None else 0 for v in defect_ppm]
    else:
        ppm = [round((1 - v / 100) * 1e6) if v is not None else 0 for v in lot_pred_yield]

    # X축 'N월 N주차': 백엔드가 대시보드 동일 앵커로 계산해 넘긴 ww_labels 우선 사용.
    # 없으면(구버전) 8/10 역산 fallback.
    if ww_labels_in and len(ww_labels_in) == n:
        ww_labels = list(ww_labels_in)
    else:
        import datetime as _dt
        _anchor_end = _dt.date(2026, 8, 10)
        ww_labels = []
        for i in range(n):
            dd = _anchor_end - _dt.timedelta(days=7 * (n - 1 - i))
            wk = (dd.day - 1) // 7 + 1
            ww_labels.append(f"{dd.month}월 {wk}주차")

    dpi = 96
    fig, ax1 = plt.subplots(figsize=(w_px/dpi, h_px/dpi), dpi=dpi)
    fig.patch.set_facecolor("white")
    ax1.set_facecolor("white")

    xs = np.arange(n)
    ax1.bar(xs, lot_production, color="#6366f1",
            alpha=0.35, width=0.55, zorder=1, label="생산량")
    ax1.set_ylabel("생산량(개)", fontsize=11, color="#94a3b8")
    ax1.tick_params(axis="y", labelsize=11, colors="#94a3b8")
    # 대시보드와 동일: y1 max 100k 고정 → 막대가 아래쪽에만 작게 표시
    ax1.set_ylim(0, 100_000)
    import matplotlib.ticker as _mticker
    ax1.yaxis.set_major_formatter(_mticker.FuncFormatter(
        lambda v, _: f"{v/1000:.0f}k" if abs(v) >= 1000 else f"{v:.0f}"
    ))

    ax2 = ax1.twinx()
    ax2.plot(xs, ppm, color="#3b82f6", linewidth=1.5, marker="o", markersize=3,
             zorder=2, label="예측불량 ppm")
    # 마지막 포인트 강조 (빨간 큰 점)
    if ppm:
        ax2.plot(xs[-1], ppm[-1], "o", color="#DC2626", markersize=7,
                 markeredgecolor="#ffffff", markeredgewidth=1.5, zorder=3)
    ax2.set_ylabel("예측불량 ppm", fontsize=11, color="#4b5563")
    ax2.tick_params(axis="y", labelsize=11, colors="#4b5563")
    # ppm y축 고정 (HTML과 동일 스케일 1500~2900)
    ax2.set_ylim(1500, 2900)
    import matplotlib.ticker as _mticker
    ax2.yaxis.set_major_formatter(_mticker.FuncFormatter(lambda v,_: f"{v:,.0f}"))
    # 각 포인트 ppm 값 라벨 (HTML ppmLabelPlugin과 동일: 예측=파랑, 최신=빨강)
    for i, v in enumerate(ppm):
        _clr = "#DC2626" if i == n - 1 else "#3b82f6"
        ax2.annotate(f"{v:,.0f}", (xs[i], v), textcoords="offset points",
                     xytext=(0, 7), ha="center", fontsize=8, fontweight="bold",
                     color=_clr, zorder=4)

    ax1.set_xticks(xs)
    ax1.set_xticklabels(ww_labels, fontsize=8, rotation=30, ha="right", color="#4b5563")
    ax1.tick_params(axis="x", length=0)
    for sp in ax1.spines.values(): sp.set_visible(False)
    for sp in ax2.spines.values(): sp.set_visible(False)
    ax1.grid(axis="y", color="#eef0f2", linewidth=0.5, zorder=0)

    lines1, labels1 = ax1.get_legend_handles_labels()
    lines2, labels2 = ax2.get_legend_handles_labels()
    ax1.legend(lines1+lines2, labels1+labels2, fontsize=11, loc="upper left",
               framealpha=0, ncol=2)

    fig.tight_layout(pad=0.3)
    buf = io.BytesIO(); fig.savefig(buf, format="png", bbox_inches="tight", dpi=dpi)
    plt.close(fig); buf.seek(0); return buf.read()


def _chart_lot_defects_png(lot_defect, w_px=580, h_px=110) -> bytes:
    """L3 Lot별 불량 개수 막대 + 불량률 꺾은선."""
    import matplotlib
    matplotlib.use("Agg")
    import matplotlib.pyplot as plt
    import numpy as np

    _try_set_font()

    if not lot_defect:
        labels = [f"Lot {40+i}" for i in range(14)]
        counts = [2,3,1,4,2,3,2,5,3,6,280,8,4,3]
        rates  = [1,2,1,2,1,2,1,3,2,3,100,4,2,2]
    else:
        labels = [d["lot"] for d in lot_defect]
        counts = [d.get("count", 0) for d in lot_defect]
        rates  = [d.get("rate",  0) for d in lot_defect]

    n = len(labels)
    max_idx = counts.index(max(counts)) if counts else 0
    colors = ["#dc2626" if i == max_idx else "rgba(148,163,184,0.5)" for i in range(n)]
    colors_mpl = ["#dc2626" if i == max_idx else "#94a3b8" for i in range(n)]

    dpi = 96
    fig, ax1 = plt.subplots(figsize=(w_px/dpi, h_px/dpi), dpi=dpi)
    fig.patch.set_facecolor("white")
    ax1.set_facecolor("white")

    xs = np.arange(n)
    ax1.bar(xs, counts, color=colors_mpl, alpha=0.8, width=0.6, zorder=1, label="Defect count")
    ax1.set_ylabel("units", fontsize=11, color="#4b5563")
    ax1.tick_params(axis="y", labelsize=11, colors="#4b5563")

    ax2 = ax1.twinx()
    ax2.plot(xs, rates, color="#111827", linewidth=1.2, marker="o", markersize=2,
             zorder=2, label="Defect rate")
    ax2.set_ylim(0, 100)
    ax2.set_ylabel("%", fontsize=11, color="#4b5563")
    ax2.tick_params(axis="y", labelsize=11, colors="#4b5563")

    ax1.set_xticks(xs)
    ax1.set_xticklabels(labels, fontsize=7, rotation=30, ha="right", color="#4b5563")
    ax1.tick_params(axis="x", length=0)
    for sp in ax1.spines.values(): sp.set_visible(False)
    for sp in ax2.spines.values(): sp.set_visible(False)
    ax1.grid(axis="y", color="#eef0f2", linewidth=0.5, zorder=0)

    lines1, labels1 = ax1.get_legend_handles_labels()
    lines2, labels2 = ax2.get_legend_handles_labels()
    ax1.legend(lines1+lines2, labels1+labels2, fontsize=11, loc="upper left",
               framealpha=0, ncol=2)

    fig.tight_layout(pad=0.3)
    buf = io.BytesIO(); fig.savefig(buf, format="png", bbox_inches="tight", dpi=dpi)
    plt.close(fig); buf.seek(0); return buf.read()


def _chart_scatter_png(feat_name, high_pts, med_pts, threshold, w_px=270, h_px=155) -> bytes:
    """L4 피처 scatter: grade1(빨강) / grade4(파랑) + 임계선."""
    import matplotlib
    matplotlib.use("Agg")
    import matplotlib.pyplot as plt

    _try_set_font()

    dpi = 96
    fig, ax = plt.subplots(figsize=(w_px/dpi, h_px/dpi), dpi=dpi)
    fig.patch.set_facecolor("white")
    ax.set_facecolor("#fafafa")

    has_data = bool(high_pts or med_pts)

    if med_pts:
        xs = [p["x"] for p in med_pts]; ys = [p["y"] for p in med_pts]
        ax.scatter(xs, ys, s=8, color="#3b82f6", alpha=0.45, label="grade4(정상)", zorder=2)
    if high_pts:
        xs = [p["x"] for p in high_pts]; ys = [p["y"] for p in high_pts]
        ax.scatter(xs, ys, s=10, color="#dc2626", alpha=0.7, label="grade1(불량)", zorder=3)
    if threshold is not None:
        ax.axvline(x=threshold, color="#dc2626", linewidth=1.2, linestyle="--", zorder=4,
                   label=f"임계값 {threshold:.4g}")

    if not has_data:
        ax.text(0.5, 0.5, "데이터 없음", ha="center", va="center",
                transform=ax.transAxes, fontsize=10, color="#9ca3af")

    ax.set_title(feat_name, fontsize=9, color="#374151", fontweight="bold", pad=3)
    ax.set_xlabel("피처값", fontsize=11, color="#4b5563", labelpad=2)
    ax.set_ylabel("reg_pred", fontsize=11, color="#4b5563", labelpad=2)
    ax.tick_params(labelsize=11, colors="#4b5563", length=2)
    for sp in ax.spines.values(): sp.set_edgecolor("#d1d5db"); sp.set_linewidth(0.5)
    ax.grid(color="#eef0f2", linewidth=0.4, zorder=0)
    if has_data:
        ax.legend(fontsize=7.5, framealpha=0.8, loc="upper right",
                  handlelength=1, borderpad=0.3, labelspacing=0.2)

    fig.tight_layout(pad=0.4)
    buf = io.BytesIO(); fig.savefig(buf, format="png", dpi=dpi, bbox_inches="tight")
    plt.close(fig); buf.seek(0); return buf.read()


def _pred_color(pred_ppm, pred_min, pred_max, threshold):
    """DrilldownV2.jsx predColor 와 동일한 연속 그라디언트 색상 반환."""
    NORMAL_STOPS = [(0.0, (243,244,246)), (0.5, (219,234,254)), (1.0, (165,215,220))]
    RISK_STOPS   = [(0.0, (254,240,138)), (0.75,(251,146, 60)), (1.0, (220, 38, 38))]

    def interp(stops, t):
        t = max(0.0, min(1.0, t))
        for i in range(1, len(stops)):
            t1, c1 = stops[i]; t0, c0 = stops[i-1]
            if t <= t1:
                k = (t - t0) / (t1 - t0 or 1)
                return tuple(int(c0[j] + (c1[j]-c0[j])*k) for j in range(3))
        return stops[-1][1]

    if pred_ppm <= threshold:
        span = max(1e-9, threshold - pred_min)
        r, g, b = interp(NORMAL_STOPS, (pred_ppm - pred_min) / span)
    else:
        span = max(1e-9, pred_max - threshold)
        r, g, b = interp(RISK_STOPS, (pred_ppm - threshold) / span)
    return f"#{r:02x}{g:02x}{b:02x}"


def _chart_custom_png(sec: dict, w_px: int, h_px: int) -> bytes:
    """custom_section dict를 matplotlib PNG로 렌더링.
    chart_type: bar / line / doughnut / pie / scatter / table.
    """
    import matplotlib
    matplotlib.use("Agg")
    import matplotlib.pyplot as plt
    import numpy as np

    _try_set_font()

    ctype     = sec.get("chart_type", "bar")
    labels    = sec.get("labels", [])
    datasets  = sec.get("datasets", [])
    horizontal = sec.get("horizontal", False)
    stacked   = sec.get("stacked", False)

    dpi = 96
    fig, ax = plt.subplots(figsize=(w_px/dpi, h_px/dpi), dpi=dpi)
    fig.patch.set_facecolor("white"); ax.set_facecolor("white")

    if ctype == "table":
        cols = sec.get("columns", [])
        rows = sec.get("rows", [])
        ax.axis("off")
        if cols and rows:
            tbl_rows = [r if isinstance(r, list) else [r.get(c, "") for c in cols] for r in rows]
            t = ax.table(cellText=tbl_rows, colLabels=cols, loc="upper left", cellLoc="left")
            t.auto_set_font_size(False); t.set_fontsize(7); t.scale(1, 1.2)
    elif ctype in ("doughnut", "pie"):
        ds0 = datasets[0] if datasets else {}
        data = ds0.get("data", [])
        colors = ds0.get("colors") or [ds0.get("color", "#3B82F6")] * len(data)
        if data:
            wedges, _ = ax.pie(data, labels=labels, colors=colors, startangle=90,
                                wedgeprops={"width": 0.4 if ctype == "doughnut" else 1.0,
                                            "edgecolor": "white", "linewidth": 1},
                                textprops={"fontsize": 7})
        ax.set_aspect("equal")
    elif ctype == "scatter":
        for ds in datasets:
            pts = ds.get("data", [])
            xs = [p.get("x", 0) for p in pts]
            ys = [p.get("y", 0) for p in pts]
            ax.scatter(xs, ys, s=10, alpha=0.6, color=ds.get("color", "#3B82F6"),
                       label=ds.get("label", ""))
        if datasets: ax.legend(fontsize=11, loc="best")
        ax.tick_params(labelsize=11, colors="#4b5563")
        ax.grid(color="#eef0f2", linewidth=0.4, zorder=0)
        for sp in ax.spines.values(): sp.set_edgecolor("#d1d5db"); sp.set_linewidth(0.5)
    else:  # bar or line
        xs = np.arange(len(labels))
        n_ds = len(datasets)
        bar_w = 0.8 / max(n_ds, 1) if not stacked else 0.7
        bottoms = np.zeros(len(labels)) if stacked else None
        for i, ds in enumerate(datasets):
            data = ds.get("data", [])
            color = ds.get("color", "#3B82F6")
            colors_ds = ds.get("colors")
            color_arg = colors_ds if (ctype == "bar" and colors_ds and len(colors_ds) == len(data)) else color
            if ctype == "line":
                ax.plot(xs, data, color=color, linewidth=1.5, marker="o", markersize=3,
                        label=ds.get("label", ""))
            else:  # bar
                if horizontal:
                    if stacked:
                        ax.barh(xs, data, left=bottoms, color=color_arg, label=ds.get("label", ""))
                        bottoms = bottoms + np.array(data)
                    else:
                        ax.barh(xs + (i - n_ds/2 + 0.5) * bar_w, data, height=bar_w,
                                color=color_arg, label=ds.get("label", ""))
                else:
                    if stacked:
                        ax.bar(xs, data, bottom=bottoms, color=color_arg, width=0.7, label=ds.get("label", ""))
                        bottoms = bottoms + np.array(data)
                    else:
                        ax.bar(xs + (i - n_ds/2 + 0.5) * bar_w, data, width=bar_w,
                               color=color_arg, label=ds.get("label", ""))
        if horizontal:
            ax.set_yticks(xs); ax.set_yticklabels(labels, fontsize=9)
            ax.invert_yaxis()
        else:
            ax.set_xticks(xs); ax.set_xticklabels(labels, fontsize=11, rotation=20, ha="right")
        ax.tick_params(labelsize=11, colors="#4b5563")
        ax.grid(axis="x" if horizontal else "y", color="#eef0f2", linewidth=0.5, zorder=0)
        ax.set_axisbelow(True)
        for sp in ax.spines.values(): sp.set_visible(False)
        if n_ds > 1 or any(ds.get("label") for ds in datasets):
            ax.legend(fontsize=11, loc="upper right", framealpha=0.85, handlelength=1.2)

    fig.tight_layout(pad=0.3)
    buf = io.BytesIO(); fig.savefig(buf, format="png", bbox_inches="tight", dpi=dpi)
    plt.close(fig); buf.seek(0); return buf.read()


def _proc_path(fname):
    import os
    return os.path.normpath(os.path.join(os.path.dirname(__file__), "..", "data", "processed", fname))


def _fi_total_gain_all():
    """feature_importance.csv 전체 X피처 lgbm_gain 합 (대시보드 ProcessFactor와 동일 분모).
    실패 시 None."""
    try:
        import pandas as pd
        df = pd.read_csv(_proc_path("feature_importance.csv"))
        df = df[df["feature"].astype(str).str.match(r"^X\d+$")]
        s = float(df["lgbm_gain"].sum())
        return s if s > 0 else None
    except Exception:
        return None


def _load_shap_bar_top(n=8):
    """shap_bar.csv에서 X피처 상위 N개 (대시보드 SHAP 영향도 기준). [{feature, mag, signed}]"""
    import pandas as pd
    df = pd.read_csv(_proc_path("shap_bar.csv"))
    df = df[df["feature"].astype(str).str.match(r"^X\d+$")]
    df = df.sort_values("mean_abs_shap", ascending=False).head(n)
    out = []
    for _, r in df.iterrows():
        out.append({"feature": str(r["feature"]),
                    "mag": float(r["mean_abs_shap"]),
                    "signed": float(r.get("mean_shap", 0) or 0)})
    return out


def _unit_red_from_beeswarm(serial, want):
    """shap_beeswarm.csv(유닛당 피처 더 많음)에서 해당 유닛의 빨강(양수) 항목.
    유닛 SHAP(저장 top10)에 빨강이 부족할 때 보충용. [{feature, mag, signed}]"""
    if not serial:
        return []
    try:
        import pandas as pd
        bee = pd.read_csv(_proc_path("shap_beeswarm.csv"),
                          usecols=["ufs_serial", "feature", "shap_value"])
        g = bee[(bee["ufs_serial"] == serial) & (bee["shap_value"] > 0)]
        return [{"feature": str(r.feature), "mag": abs(float(r.shap_value)),
                 "signed": float(r.shap_value)} for r in g.itertuples()]
    except Exception:
        return []


def _shap_red_items(report_data, n):
    """SHAP 영향도 빨강(양수 SHAP=불량↑)만 필터 → mag 내림차순 → 상위 n개.
    유닛 SHAP(저장 top10)에 빨강이 n개 미만이면 beeswarm에서 빨강을 보충."""
    _us = report_data.get("unit_shap")
    if _us:
        reds = [dict(it) for it in _us if it.get("signed", 0) >= 0]
        if len(reds) < n:
            _bred = _unit_red_from_beeswarm((report_data.get("top_unit") or {}).get("serial"), n)
            if len(_bred) > len(reds):
                reds = _bred
    else:
        reds = [it for it in _load_shap_bar_top(max(n * 4, 30)) if it.get("signed", 0) >= 0]
    reds.sort(key=lambda it: it.get("mag", 0), reverse=True)   # 영향도 큰 순 정렬
    return reds[:n]


def _top_red_shap_feature(report_data):
    """SHAP 영향도 빨강(양수) 중 1등 피처명 (분포차트 라벨용)."""
    reds = _shap_red_items(report_data, 1)
    return reds[0].get("feature", "") if reds else ""


def _dashboard_feat_dist(feature=None, bins=40):
    """대시보드 '안전 vs 위험 분포'와 동일 — 최상위 SHAP 피처의
    위험(reg_pred 상위10%) vs 안전(하위10%) 값 분포 + 교차점 임계값."""
    import pandas as pd, numpy as np
    fd = pd.read_csv(_proc_path("feature_dist.csv"))
    u = pd.read_csv(_proc_path("dashboard_units.csv"), usecols=["ufs_serial", "reg_pred"])
    u["reg_pred"] = pd.to_numeric(u["reg_pred"], errors="coerce")
    reg = u.dropna(subset=["reg_pred"]).set_index("ufs_serial")["reg_pred"]
    if reg.empty:
        return None
    p90, p10 = float(reg.quantile(0.90)), float(reg.quantile(0.10))
    if feature is None:
        sb = pd.read_csv(_proc_path("shap_bar.csv"))
        sb = sb[sb["feature"].astype(str).str.match(r"^X\d+$")].sort_values("mean_abs_shap", ascending=False)
        for f in sb["feature"]:
            if f in fd.columns:
                feature = f
                break
    if feature is None or feature not in fd.columns:
        return None
    m = fd[["ufs_serial", feature]].copy()
    m["rp"] = m["ufs_serial"].map(reg)
    m = m.dropna(subset=[feature, "rp"])
    hi = m[m["rp"] >= p90][feature].values
    lo = m[m["rp"] <= p10][feature].values
    if len(hi) < 10 or len(lo) < 10:
        return None
    allv = np.concatenate([hi, lo])
    mn, mx = np.quantile(allv, 0.01), np.quantile(allv, 0.99)
    if not (mx > mn):
        return None
    edges = np.linspace(mn, mx, bins + 1)
    centers = ((edges[:-1] + edges[1:]) / 2)

    def dens(v):
        idx = np.clip(np.digitize(v, edges) - 1, 0, bins - 1)
        c = np.bincount(idx, minlength=bins).astype(float)
        return c / max(len(v), 1) * 100

    dHi, dLo = dens(hi), dens(lo)
    hiMed, loMed = float(np.median(hi)), float(np.median(lo))
    direction = "up" if hiMed > loMed else "down"
    c0, c1 = min(loMed, hiMed), max(loMed, hiMed)
    thr = hiMed if direction == "up" else loMed
    for i in range(1, bins):
        xc = centers[i]
        if xc < c0 or xc > c1:
            continue
        prev = dHi[i - 1] - dLo[i - 1]
        cur = dHi[i] - dLo[i]
        if direction == "up" and prev < 0 and cur >= 0:
            thr = xc; break
        if direction == "down" and prev >= 0 and cur < 0:
            thr = xc; break
    return {"feature": feature, "labels": centers.tolist(),
            "danger": dHi.tolist(), "normal": dLo.tolist(), "threshold": float(thr)}


def _chart_shap_bar_png(items, w_px=300, h_px=240) -> bytes:
    """대시보드 'SHAP 영향도' 핵심 — 피처별 평균 |SHAP| 수평 바 (ppm), 방향=색."""
    import matplotlib
    matplotlib.use("Agg")
    import matplotlib.pyplot as plt
    import numpy as np

    _try_set_font()
    dpi = 96
    fig, ax = plt.subplots(figsize=(w_px / dpi, h_px / dpi), dpi=dpi)
    fig.patch.set_facecolor("white"); ax.set_facecolor("white")

    if not items:
        ax.text(0.5, 0.5, "데이터 없음", ha="center", va="center",
                transform=ax.transAxes, fontsize=10, color="#9ca3af")
    else:
        labels = [it["feature"] for it in items]
        vals = [it["mag"] * 1e6 for it in items]            # ppm
        colors = ["#EF4444" if it.get("signed", 0) >= 0 else "#3B82F6" for it in items]
        ys = np.arange(len(labels))
        ax.barh(ys, vals, color=colors, height=0.6)
        ax.set_yticks(ys)
        ax.set_yticklabels(labels, fontsize=10, color="#111827", fontweight="bold")
        ax.invert_yaxis()
        ax.set_xlim(0, (max(vals) if vals else 1) * 1.15)
        ax.set_xlabel("평균 |SHAP| (ppm)", fontsize=10.5, color="#4b5563", labelpad=2)
        ax.tick_params(axis="x", labelsize=10, colors="#6b7280")
        ax.tick_params(axis="y", length=0)
        for sp in ax.spines.values(): sp.set_visible(False)
        ax.grid(axis="x", color="#eef0f2", linewidth=0.5, zorder=0)
        ax.set_axisbelow(True)

    fig.tight_layout(pad=0.3)
    buf = io.BytesIO(); fig.savefig(buf, format="png", bbox_inches="tight", dpi=dpi)
    plt.close(fig); buf.seek(0); return buf.read()


def _chart_fi_bar_png(features, w_px=580, h_px=200) -> bytes:
    """L2 Feature Importance Top N 수평 바 차트 (HTML c-fi-top와 동일)."""
    import matplotlib
    matplotlib.use("Agg")
    import matplotlib.pyplot as plt
    import numpy as np

    _try_set_font()

    feats_top = features[:5]
    total_all = _fi_total_gain_all() or sum(f.get("lgbm_gain", 0) or 0 for f in features) or 1
    labels = [f.get("feature", "") for f in feats_top]
    values = [round((f.get("lgbm_gain", 0) or 0) / total_all * 100, 2) for f in feats_top]
    max_v  = max(values) if values else 1

    dpi = 96
    fig, ax = plt.subplots(figsize=(w_px/dpi, h_px/dpi), dpi=dpi)
    fig.patch.set_facecolor("white"); ax.set_facecolor("white")

    if not labels:
        ax.text(0.5, 0.5, "데이터 없음", ha="center", va="center",
                transform=ax.transAxes, fontsize=10, color="#9ca3af")
    else:
        ys = np.arange(len(labels))
        colors = ["#1e3a5f" if i == 0 else "#374151" for i in range(len(labels))]
        ax.barh(ys, values, color=colors, height=0.55)
        ax.set_yticks(ys)
        ax.set_yticklabels(labels, fontsize=10, color="#111827", fontweight="bold")
        ax.invert_yaxis()
        ax.set_xlim(0, max_v * 1.25)
        from matplotlib.ticker import FuncFormatter
        ax.xaxis.set_major_formatter(FuncFormatter(lambda v, _: f"{v:.1f}%"))
        ax.tick_params(axis="x", labelsize=12, colors="#6b7280")
        ax.tick_params(axis="y", length=0)
        for sp in ax.spines.values(): sp.set_visible(False)
        ax.grid(axis="x", color="#eef0f2", linewidth=0.5, zorder=0)
        # 막대 옆 % 라벨 (HTML fiLabelPlugin과 동일)
        for yi, v in zip(ys, values):
            ax.text(v + max_v * 0.02, yi, f"{v:.1f}%", va="center", ha="left",
                    fontsize=10, fontweight="bold", color="#374151")
        ax.set_axisbelow(True)

    fig.tight_layout(pad=0.3)
    buf = io.BytesIO(); fig.savefig(buf, format="png", bbox_inches="tight", dpi=dpi)
    plt.close(fig); buf.seek(0); return buf.read()


def _chart_fdc_line_png(labels, normal, danger, threshold, feature_name, w_px=300, h_px=240) -> bytes:
    """R3 피처 정상/불량 분포 비교 라인 차트 (HTML c-feat-dist와 동일)."""
    import matplotlib
    matplotlib.use("Agg")
    import matplotlib.pyplot as plt
    import numpy as np

    _try_set_font()

    dpi = 96
    fig, ax = plt.subplots(figsize=(w_px/dpi, h_px/dpi), dpi=dpi)
    fig.patch.set_facecolor("white"); ax.set_facecolor("white")

    if not labels:
        ax.text(0.5, 0.5, "데이터 없음", ha="center", va="center",
                transform=ax.transAxes, fontsize=10, color="#9ca3af")
    else:
        xs = np.array(labels, dtype=float)
        ax.fill_between(xs, normal, color="#3B82F6", alpha=0.12)
        ax.plot(xs, normal, color="#3B82F6", linewidth=1.5, label="정상")
        ax.fill_between(xs, danger, color="#EF4444", alpha=0.12)
        ax.plot(xs, danger, color="#EF4444", linewidth=1.5, label="위험")
        if threshold is not None:
            ax.axvline(x=float(threshold), color="#dc2626", linewidth=1.2, linestyle="--")
            ax.text(float(threshold), ax.get_ylim()[1]*0.95,
                    f" 임계 {float(threshold):.2f}",
                    color="#dc2626", fontsize=11.5, fontweight="bold", va="top")
        ax.set_xlabel("피처값", fontsize=11.5, color="#4b5563", labelpad=2)
        ax.set_ylabel("비율 (%)", fontsize=11.5, color="#4b5563", labelpad=2)
        ax.tick_params(labelsize=11, colors="#4b5563", length=2)
        # 범례를 차트 위쪽 바깥에 배치 (차트 데이터 가리지 않도록)
        ax.legend(fontsize=9, framealpha=0, loc="lower center",
                  bbox_to_anchor=(0.5, 1.02), ncol=2,
                  handlelength=1.2, borderpad=0.2, labelspacing=0.2,
                  columnspacing=1.0)
        for sp in ax.spines.values(): sp.set_edgecolor("#d1d5db"); sp.set_linewidth(0.5)
        ax.grid(color="#eef0f2", linewidth=0.4, zorder=0)

    fig.tight_layout(pad=0.3)
    buf = io.BytesIO(); fig.savefig(buf, format="png", bbox_inches="tight", dpi=dpi)
    plt.close(fig); buf.seek(0); return buf.read()


def _chart_wafer_png(wd_dies, wd_x_range, wd_y_range, serial, w_px=240, h_px=240) -> bytes:
    """웨이퍼맵: HTML SVG 렌더링과 동일한 좌표·스타일 (정사각형 출력, 컬러바 없음).
    HTML build_html의 wafer_svg와 1:1 매핑 — 픽셀 좌표계, y↓.
    """
    import matplotlib
    matplotlib.use("Agg")
    import matplotlib.pyplot as plt
    from matplotlib.patches import Circle, Rectangle
    from matplotlib.patches import PathPatch
    from matplotlib.path import Path as MPath

    _try_set_font()

    # pred 색상 스케일 — wafer_scale.json 글로벌 기준
    preds = [d.get("pred") if d.get("pred") is not None else (d.get("pred_ppm", 0) or 0) / 1e6
             for d in wd_dies]
    try:
        import json as _jmod, os as _omod
        _sp = _omod.path.normpath(_omod.path.join(_omod.path.dirname(__file__), "..", "data", "processed", "wafer_scale.json"))
        with open(_sp) as _f:
            _ws = _jmod.load(_f)
        pred_min  = float(_ws.get("pred_min", 0))
        pred_max  = float(_ws.get("pred_max", 0.006))
        threshold = float(_ws.get("threshold", 0.003))
    except Exception:
        pred_min = min(preds) if preds else 0
        pred_max = max(preds) if preds else 1
        sorted_p = sorted(preds)
        threshold = sorted_p[int(len(sorted_p) * 0.75)] if sorted_p else (pred_max * 0.5)

    # HTML SVG와 완전 동일한 좌표 계산
    side = min(w_px, h_px)
    wmap_w = wmap_h = side
    cx_px, cy_px = wmap_w / 2, wmap_h / 2
    pad = 14
    r_svg = min(cx_px, cy_px) - pad

    # GLOBAL die grid → pixel 변환 (DrilldownV2 SCALE=0.9)
    GX_MIN, GX_MAX = 12, 66
    GY_MIN, GY_MAX = 11, 32
    ref_x = GX_MAX - GX_MIN + 1   # 55
    ref_y = GY_MAX - GY_MIN + 1   # 22
    mx = (GX_MIN + GX_MAX) / 2    # 39
    my = (GY_MIN + GY_MAX) / 2    # 21.5
    SCALE = 0.9
    sx = (r_svg * 2 / ref_x) * SCALE
    sy = (r_svg * 2 / ref_y) * SCALE
    die_w = sx
    die_h = sy

    def to_px(dx, dy):
        return cx_px + (dx - mx) * sx, cy_px + (dy - my) * sy

    # figure: 정사각형, 축은 figure 전체 차지, 마진 0
    dpi = 96
    fig, ax = plt.subplots(figsize=(wmap_w/dpi, wmap_h/dpi), dpi=dpi)
    fig.patch.set_facecolor("white")
    ax.set_facecolor("white")
    fig.subplots_adjust(left=0, right=1, top=1, bottom=0)
    ax.set_xlim(0, wmap_w)
    ax.set_ylim(wmap_h, 0)  # y↓ (SVG 좌표)
    ax.set_aspect("equal")
    ax.axis("off")

    # 웨이퍼 배경 원
    wafer_bg = Circle((cx_px, cy_px), r_svg, fill=True,
                      facecolor="#F8FAFC", edgecolor="#94A3B8", linewidth=1.5)
    ax.add_patch(wafer_bg)

    # 원 클립 영역 (SVG clipPath와 동일)
    clip_path = mpath_circle = MPath.unit_circle()
    from matplotlib.transforms import Affine2D
    clip_transform = Affine2D().scale(r_svg, r_svg).translate(cx_px, cy_px) + ax.transData
    clip_patch = PathPatch(clip_path, transform=clip_transform, facecolor='none', edgecolor='none')
    ax.add_patch(clip_patch)

    # die 사각형 (clip 적용)
    for die in wd_dies:
        px, py = to_px(die["x"], die["y"])
        die_pred = die.get("pred") if die.get("pred") is not None else (die.get("pred_ppm", 0) or 0) / 1e6
        clr = _pred_color(die_pred, pred_min, pred_max, threshold)
        is_target = die.get("is_target", False)
        rect = Rectangle((px - die_w/2, py - die_h/2), die_w, die_h,
                          facecolor=clr,
                          edgecolor="#7f1d1d" if is_target else "none",
                          linewidth=0.8 if is_target else 0)
        rect.set_clip_path(clip_patch)
        ax.add_patch(rect)

    # 노치 (웨이퍼 하단 작은 사각형)
    notch = Rectangle((cx_px - 3, cy_px + r_svg - 3), 6, 6,
                       facecolor="#CBD5E1", edgecolor="none")
    ax.add_patch(notch)

    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=dpi, pad_inches=0)
    plt.close(fig); buf.seek(0); return buf.read()


def _chart_wafer_cbar_png(w_px=24, h_px=200) -> bytes:
    """웨이퍼맵 옆 세로 컬러바 (정상→위험). 좁은 그라디언트 + 위아래 라벨."""
    import matplotlib
    matplotlib.use("Agg")
    import matplotlib.pyplot as plt
    from matplotlib.colors import LinearSegmentedColormap
    import numpy as np

    _try_set_font()
    dpi = 96
    fig, ax = plt.subplots(figsize=(w_px/dpi, h_px/dpi), dpi=dpi)
    fig.patch.set_facecolor("white")
    # 라벨 공간 확보: 좌우 여유 두고 그라디언트는 가운데 좁게
    fig.subplots_adjust(left=0.30, right=0.70, top=0.88, bottom=0.12)

    grad_colors = ["#f3f4f6", "#dbeafe", "#a5d7dc", "#fef08a", "#fb923c", "#dc2626"]
    cmap = LinearSegmentedColormap.from_list("wafer", grad_colors)
    gradient = np.linspace(1, 0, 256).reshape(256, 1)
    ax.imshow(gradient, aspect="auto", cmap=cmap)
    ax.set_xticks([]); ax.set_yticks([])
    for sp in ax.spines.values():
        sp.set_edgecolor("#94a3b8"); sp.set_linewidth(0.5)

    # 위/아래 라벨 (figure 전체 좌표 기준으로 가운데 정렬)
    fig.text(0.5, 0.96, "위험", ha="center", va="top",
             fontsize=7, fontweight="bold", color="#dc2626")
    fig.text(0.5, 0.04, "정상", ha="center", va="bottom",
             fontsize=7, fontweight="bold", color="#16803c")

    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=dpi, pad_inches=0)
    plt.close(fig); buf.seek(0); return buf.read()


_FONT_READY = False

def _try_set_font():
    """matplotlib 한글 폰트 설정.
    배포 서버(Linux)에 한글 폰트가 없어 범례가 □로 깨지는 문제 →
    repo 번들 폰트(agent/fonts/NanumGothic.ttf)를 직접 등록하여 플랫폼 무관하게 한글 보장."""
    global _FONT_READY
    import os
    import matplotlib.pyplot as plt
    import matplotlib.font_manager as fm
    import platform

    plt.rcParams["axes.unicode_minus"] = False

    # 1) repo 번들 폰트 (배포 Linux 등 시스템 한글 폰트 없어도 동작)
    bundled = os.path.join(os.path.dirname(__file__), "fonts", "NanumGothic.ttf")
    if os.path.exists(bundled):
        try:
            if not _FONT_READY:
                fm.fontManager.addfont(bundled)
                _FONT_READY = True
            plt.rcParams["font.family"] = fm.FontProperties(fname=bundled).get_name()
            return
        except Exception:
            pass

    # 2) 시스템 폰트 폴백
    candidates = (
        ["Malgun Gothic"] if platform.system() == "Windows"
        else ["NanumGothic", "AppleGothic"]
    )
    for fn in candidates:
        try:
            fm.findfont(fm.FontProperties(family=fn), fallback_to_default=False)
            plt.rcParams["font.family"] = fn
            return
        except Exception:
            pass
    plt.rcParams["font.family"] = "DejaVu Sans"


# ── PPTX 헬퍼 ─────────────────────────────────────────────────
def _emu(px: float, dpi: float = 96.0) -> int:
    """픽셀 → EMU (1 inch = 914400 EMU, 1 inch = 96px)."""
    return int(px / dpi * 914400)


def _png_shape(slide, png_bytes: bytes, x_px, y_px, w_px, h_px, dpi=96):
    """PNG bytes를 PPTX 슬라이드에 삽입."""
    from pptx.util import Emu
    from io import BytesIO
    slide.shapes.add_picture(
        BytesIO(png_bytes),
        Emu(_emu(x_px, dpi)), Emu(_emu(y_px, dpi)),
        Emu(_emu(w_px, dpi)), Emu(_emu(h_px, dpi)),
    )


def _box(slide, x_px, y_px, w_px, h_px, fill_rgb, line_rgb=None, line_pt=0.5, radius=False, dpi=96):
    from pptx.util import Emu, Pt as _Pt
    from pptx.dml.color import RGBColor as RGB
    shape = slide.shapes.add_shape(
        1 if not radius else 5,
        Emu(_emu(x_px, dpi)), Emu(_emu(y_px, dpi)),
        Emu(_emu(w_px, dpi)), Emu(_emu(h_px, dpi))
    )
    shape.fill.solid(); shape.fill.fore_color.rgb = RGB(*fill_rgb)
    if line_rgb:
        shape.line.color.rgb = RGB(*line_rgb); shape.line.width = _Pt(line_pt)
    else:
        shape.line.fill.background()
    return shape


def _txt(slide, text, x_px, y_px, w_px, h_px, size=9, bold=False, color=(32,36,42),
         align="left", wrap=True, dpi=96):
    from pptx.util import Emu, Pt as _Pt
    from pptx.dml.color import RGBColor as RGB
    from pptx.enum.text import PP_ALIGN as _ALIGN
    tb = slide.shapes.add_textbox(
        Emu(_emu(x_px, dpi)), Emu(_emu(y_px, dpi)),
        Emu(_emu(w_px, dpi)), Emu(_emu(h_px, dpi))
    )
    tf = tb.text_frame; tf.word_wrap = wrap
    # 텍스트 프레임 기본 내부 마진 제거 (텍스트가 박스 좌상단에 정확히 위치)
    tf.margin_top = 0; tf.margin_bottom = 0
    tf.margin_left = 0; tf.margin_right = 0
    p = tf.paragraphs[0]
    p.alignment = {"left":_ALIGN.LEFT,"center":_ALIGN.CENTER,"right":_ALIGN.RIGHT}.get(align, _ALIGN.LEFT)
    run = p.add_run(); run.text = _strip_html(text)
    run.font.size = _Pt(size); run.font.bold = bold; run.font.color.rgb = RGB(*color)
    # remove default Arial fallback
    try: run.font.name = "Malgun Gothic"
    except Exception: pass
    return tb


def _strip_html(s: str) -> str:
    import re
    s = re.sub(r"<[^>]+>", "", s)
    s = s.replace("&nbsp;", " ").replace("&amp;", "&").replace("&lt;", "<").replace("&gt;", ">")
    return s


def _multi_run_tx(slide, runs, x_px, y_px, w_px, h_px, sz=10, bold=False,
                  align="left", wrap=True, dpi=96):
    """여러 색상 run을 한 텍스트 박스에 표시. runs=[(text, (r,g,b)), ...]."""
    from pptx.util import Emu, Pt as _Pt
    from pptx.dml.color import RGBColor as RGB
    from pptx.enum.text import PP_ALIGN as _ALIGN
    tb = slide.shapes.add_textbox(
        Emu(_emu(x_px, dpi)), Emu(_emu(y_px, dpi)),
        Emu(_emu(w_px, dpi)), Emu(_emu(h_px, dpi))
    )
    tf = tb.text_frame; tf.word_wrap = wrap
    tf.margin_top = 0; tf.margin_bottom = 0
    tf.margin_left = 0; tf.margin_right = 0
    p = tf.paragraphs[0]
    p.alignment = {"left":_ALIGN.LEFT,"center":_ALIGN.CENTER,"right":_ALIGN.RIGHT}.get(align, _ALIGN.LEFT)
    for text, color in runs:
        if not text:
            continue
        run = p.add_run()
        run.text = _strip_html(text)
        run.font.size = _Pt(sz); run.font.bold = bold
        run.font.color.rgb = RGB(*color)
        try: run.font.name = "Malgun Gothic"
        except Exception: pass
    return tb


def _boxtxt(slide, text, x_px, y_px, w_px, h_px, fill_rgb, text_rgb=(255,255,255),
            size=9, bold=False, align="left", dpi=96):
    """배경색 박스 + 텍스트."""
    _box(slide, x_px, y_px, w_px, h_px, fill_rgb, dpi=dpi)
    _txt(slide, text, x_px+3, y_px+1, w_px-6, h_px-2, size=size, bold=bold,
         color=text_rgb, align=align, dpi=dpi)


# ── 메인 빌드 ─────────────────────────────────────────────────
def _apply_html_overrides(report_data: dict, html: str) -> dict:
    """현재 미리보기 HTML에서 인라인 편집된 텍스트들을 추출해 report_data 덮어쓰기.
    인라인 텍스트 편집(텍스트 더블클릭으로 직접 수정)이 PPT에도 반영되도록 함.
    """
    import re as _re_h, copy as _copy_h
    d = _copy_h.deepcopy(report_data)
    d.setdefault("meta", {})

    def _extract_inner(pattern):
        m = _re_h.search(pattern, html, _re_h.DOTALL)
        if not m: return None
        # HTML 태그 제거
        inner = _re_h.sub(r"<[^>]+>", "", m.group(1)).strip()
        inner = inner.replace("&nbsp;", " ").replace("&amp;", "&")
        inner = inner.replace("&lt;", "<").replace("&gt;", ">")
        return inner if inner else None

    # 헤더 제목
    v = _extract_inner(r'<div class="s-caption"[^>]*>(.*?)</div>')
    if v: d["meta"]["report_title"] = v

    # 배너
    v = _extract_inner(r'<div class="s-title"[^>]*>(.*?)</div>')
    if v: d["meta"]["summary_title"] = v
    v = _extract_inner(r'<div class="s-subtitle"[^>]*>(.*?)</div>')
    if v: d["meta"]["summary_sub"] = v

    # 섹션 헤더 [모델링 결과] / [불량 예측 현황]
    for label_key, pat in [
        ("left_header",  r'<div class="shdr"[^>]*>(\[ ?모델링[^<]*)'),
        ("right_header", r'<div class="shdr"[^>]*>(\[ ?불량[^<]*)'),
    ]:
        v = _extract_inner(pat + r'(?:.*?)</div>')
        if v:
            d["meta"].setdefault("section_labels", {})[label_key] = v

    # 좌측 inum 섹션 라벨 (2. xxx, 3. xxx)
    inums = _re_h.findall(r'<div class="inum"[^>]*>(.*?)</div>', html, _re_h.DOTALL)
    for inum_html in inums:
        plain = _re_h.sub(r"<[^>]+>", "", inum_html).strip()
        m = _re_h.match(r'(\d+)\.\s*(.+?)(?:\s+출처:.*)?$', plain)
        if not m: continue
        num, label = m.group(1), m.group(2).strip()
        sl_map = {"1": "L1", "2": "L2", "3": "L3"}
        if num in sl_map:
            d["meta"].setdefault("section_labels", {})[sl_map[num]] = label

    return d


def build_pptx(report_data: dict, current_html: str | None = None) -> bytes:
    """
    HTML 보고서와 1:1 동일 레이아웃 PPTX 생성.
    차트는 matplotlib PNG, 텍스트/박스/테이블은 pptx 도형으로 재현.
    current_html: 미리보기 iframe의 현재 HTML (인라인 텍스트 편집 반영용)
    """
    from datetime import timedelta
    import re

    # 인라인 편집 내용을 report_data에 덮어쓰기 (PPT에 반영)
    if current_html:
        try:
            report_data = _apply_html_overrides(report_data, current_html)
        except Exception as _e:
            print(f"[build_pptx] HTML override 파싱 실패: {_e}")

    meta         = report_data.get("meta", {})
    scan         = report_data.get("scan", {})
    importance   = report_data.get("importance", {})
    analysis     = report_data.get("analysis", {})
    features     = importance.get("features", [])
    top_features = analysis.get("top_features", [])

    val_rmse  = meta.get("val_rmse",  "0.005699")
    model_nm  = meta.get("model",     "Stacking Ensemble")
    today     = datetime(2026, 6, 11)   # 오늘 고정
    today_str = today.strftime("%Y. %m. %d")
    scan_total = scan.get("total_units", "-")
    scan_high  = scan.get("grade1_count", scan.get("high_count", "-"))

    ppm_delta     = report_data.get("ppm_delta", {})
    delta_val     = ppm_delta.get("delta", 0)
    delta_str     = f"▲{abs(delta_val):,.0f}" if delta_val >= 0 else f"▼{abs(delta_val):,.0f}"
    try:
        _default_alert = ", ".join(it["feature"] for it in _shap_red_items(report_data, 1))   # SHAP 차트 상위 1개(빨강=불량↑)
    except Exception:
        _default_alert = ", ".join(ppm_delta.get("top_features", [f.get("feature","") for f in top_features[:2]])) or "-"
    alert_features = meta.get("alert_features", _default_alert)
    report_title   = meta.get("report_title", "Field Health 불량 예측 분석 보고서")
    summary_title  = _strip_html(meta.get("summary_title",
        f"전 주 대비 품질 불량 {delta_str} ppm {'열화' if delta_val>=0 else '개선'}"))
    summary_sub    = _strip_html(meta.get("summary_sub",
        f"원인 WT Parameter {alert_features} 이상 → inline 참원인 도출 요청"))

    _sl     = meta.get("section_labels", {})
    slabel  = lambda key, default: _sl.get(key, default)

    # PPM 계산 — 전체 평균 reg_pred를 ppm으로 환산
    try:
        from tools import get_mean_pred_ppm
        _ppm_val = get_mean_pred_ppm()
        _ppm_str = f"{_ppm_val:,}" if _ppm_val > 0 else "-"
    except Exception:
        _ppm_str = "-"

    # ── 데이터 준비
    weekly_yield  = report_data.get("weekly_yield_trend", {})
    lot_labels    = weekly_yield.get("labels",     [])
    lot_production= weekly_yield.get("production", [])
    lot_pred_yield= weekly_yield.get("pred_yield", [])
    lot_defect_ppm= weekly_yield.get("defect_ppm", [])
    lot_ww_labels = weekly_yield.get("ww_labels",  [])   # 'N월 N주차' (대시보드 앵커)
    lot_date_labels = weekly_yield.get("date_labels", [])

    lot_defect    = report_data.get("lot_defect_counts", [])
    if not lot_defect:
        pred_ppm_trend = report_data.get("pred_ppm_trend", {})
        _lbals = pred_ppm_trend.get("labels", [f"Lot {40+i}" for i in range(14)])
        _lhigh = pred_ppm_trend.get("defect_count", pred_ppm_trend.get("high_ppm", []))
        _lrate = pred_ppm_trend.get("defect_rate", [])
        _lhigh_safe = (_lhigh or []) + [0] * max(0, len(_lbals) - len(_lhigh or []))
        _lrate_safe = (_lrate or []) + [0] * max(0, len(_lbals) - len(_lrate or []))
        lot_defect = [{"lot": l, "count": c, "rate": r}
                      for l, c, r in zip(_lbals, _lhigh_safe, _lrate_safe)]

    feat_scatter = report_data.get("feat_scatter", {})
    fs1 = feat_scatter.get("feat1", {}); fs2 = feat_scatter.get("feat2", {})
    fs1_name = fs1.get("name","Feat1"); fs2_name = fs2.get("name","Feat2")
    fs1_threshold = fs1.get("threshold"); fs2_threshold = fs2.get("threshold")
    fs1_high = fs1.get("pts_high",[]); fs1_med = fs1.get("pts_med",[])
    fs2_high = fs2.get("pts_high",[]); fs2_med = fs2.get("pts_med",[])

    wafer_die  = report_data.get("wafer_die", {})
    wd_dies    = wafer_die.get("dies", [])
    wd_x_range = wafer_die.get("x_range", [12,66])
    wd_y_range = wafer_die.get("y_range", [11,32])

    _tu          = report_data.get("top_unit", {})
    _pred_h_raw  = _tu.get("pred_health", _tu.get("pred_ppm", 0))
    try: _pred_h_f = float(_pred_h_raw)
    except: _pred_h_f = 0.0
    _pred_h_frac = (_pred_h_f / 1e6) if _pred_h_f > 1 else _pred_h_f
    _pred_h_disp = f"{round(_pred_h_frac * 1e6):,} ppm"
    unit_serial  = _tu.get("serial", "S38369")
    unit_lot     = str(_tu.get("run_id", "-"))
    unit_wafer   = str(_tu.get("wafer_no", "-"))

    anomaly_stats = report_data.get("anomaly_stats", [])
    if not anomaly_stats:
        _top_feat_map = {f.get("feature",""): f for f in top_features}
        for f in features[:3]:
            fname = f.get("feature","")
            _af   = _top_feat_map.get(fname, {})
            ratio = _af.get("ratio")
            if ratio is None:
                gain = f.get("lgbm_gain", 0) or 0
                tg   = sum(x.get("lgbm_gain",0) or 0 for x in features) or 1
                danger = min(int(gain/tg*1000), 90)
            else:
                danger = min(int(abs(ratio-1.0)/3.0*100), 95)
            anomaly_stats.append({"feature": fname, "danger": danger, "normal": 100-danger})

    total_gain = _fi_total_gain_all() or sum(f.get("lgbm_gain",0) or 0 for f in features) or 1
    max_gain   = max((f.get("lgbm_gain",0) or 0 for f in features), default=1) or 1

    # ── 보고서 수정 상태 (hidden / custom / commentary) ─────────
    hidden_sids     = set(report_data.get("hidden_sections", []) or [])
    custom_sections = report_data.get("custom_sections", []) or []
    commentary_list = report_data.get("commentary", []) or []
    replace_map = {s["replace_sid"]: s for s in custom_sections if s.get("replace_sid")}
    extra_sections = [s for s in custom_sections if not s.get("replace_sid")]

    # ── 슬라이드 좌표 상수 ─────────────────────────────────────
    # 슬라이드: 1280×720px  DPI=96
    prs   = _new_prs()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    DPI   = 96.0
    SW, SH = 1280, 720

    def bx(x, y, w, h, fill, line=None, lpt=0.5):
        return _box(slide, x, y, w, h, fill, line, lpt, dpi=DPI)
    def tx(text, x, y, w, h, sz=11, bold=False, clr=(32,36,42), align="left", wrap=True):
        return _txt(slide, text, x, y, w, h, sz, bold, clr, align, wrap, DPI)
    def btx(text, x, y, w, h, fill, tclr=(255,255,255), sz=11, bold=False, align="left"):
        return _boxtxt(slide, text, x, y, w, h, fill, tclr, sz, bold, align, DPI)
    def img(png_bytes, x, y, w, h):
        return _png_shape(slide, png_bytes, x, y, w, h, DPI)

    # ─── 전역 좌표 ──────────────────────────────────────────────
    TOPBAR_H = 36
    SUM_H    = 50
    BODY_Y   = TOPBAR_H + SUM_H      # 86
    FTR_H    = 22
    BODY_H   = SH - BODY_Y - FTR_H  # 612
    PAD      = 10
    GAP      = 10
    COL_W    = (SW - PAD*2 - GAP) // 2   # 625
    LX       = PAD                        # 10
    RX       = PAD + COL_W + GAP          # 645
    SBOX_H   = BODY_H - 4                 # 608

    # ─── 헤더 ──────────────────────────────────────────────────
    bx(0, 0, SW, SH, (255,255,255))
    bx(0, 0, SW, TOPBAR_H, (255,255,255))
    bx(0, TOPBAR_H-2, SW, 2, (30,58,138))
    tx(f"발행일자: {today_str}", 18, 10, 220, 18, sz=11, bold=True, clr=(51,78,118))
    tx(report_title, 0, 6, SW, 26, sz=15, bold=True, clr=(15,23,42), align="center")
    btx("대외비", SW-100, 8, 86, 22, fill=(220,38,38), sz=11, bold=True, align="center")

    # ─── 요약 배너 ─────────────────────────────────────────────
    bx(0, TOPBAR_H, SW, SUM_H, (255,247,237))
    bx(0, TOPBAR_H+SUM_H-1, SW, 1, (253,215,170))
    # summary_title: "▲N ppm" 부분만 색상 강조 (▲=빨강 열화, ▼=초록 개선)
    import re as _re_st
    _m = _re_st.search(r'([▲▼]\s*\d[\d,]*\s*ppm)', summary_title)
    if _m:
        _pre = summary_title[:_m.start()]
        _mid = _m.group(1)
        _suf = summary_title[_m.end():]
        _delta_clr = (220,38,38) if _mid.startswith('▲') else (22,163,74)
        _runs = [(_pre, (17,24,39)), (_mid, _delta_clr), (_suf, (17,24,39))]
        _multi_run_tx(slide, _runs, 0, TOPBAR_H+5, SW, 24, sz=15, bold=True, align="center", dpi=DPI)
    else:
        tx(summary_title, 0, TOPBAR_H+5, SW, 24, sz=15, bold=True, clr=(17,24,39), align="center")
    tx(summary_sub,   0, TOPBAR_H+30, SW, 18, sz=11, bold=True, clr=(26,58,92), align="center")

    # ─── 좌칼럼 sbox ───────────────────────────────────────────
    LW = COL_W
    BY = BODY_Y + 2

    bx(LX, BY, LW, SBOX_H, (255,255,255), (107,114,128), 0.5)
    SHDR_H = 30
    bx(LX, BY, LW, SHDR_H, (243,244,246), (107,114,128), 0.5)
    bx(LX, BY, 4, SHDR_H, (107,114,128))
    tx(slabel("left_header","[ 모델링 결과 ]"), LX+12, BY+8, LW-24, 16, sz=12, bold=True, clr=(17,24,39))

    # 좌칼럼 내부 y 커서 (패딩 10px)
    cy = BY + SHDR_H + 10

    # ── 1. 모델 성능 KPI ──────────────────────────────────────
    tx(f"{slabel('L1','모델 성능')}", LX+12, cy, LW-24, 16, sz=11, bold=True, clr=(17,24,39))
    cy += 18

    KPI_H = 56
    KPI_GAP = 8
    KPI_W = (LW - 24 - KPI_GAP*2) // 3
    # HTML과 동일: sub 텍스트 제거
    kpi_data = [
        ("RMSE",            val_rmse,           (30,58,138), (30,58,138)),
        ("분석 유닛",       f"{scan_total}개",  (55,65,81),  (55,65,81)),
        ("평균 예측 health", _ppm_str + " ppm", (180,83,9),  (180,83,9)),
    ]
    for i, (lbl, val, lclr, vclr) in enumerate(kpi_data):
        kx = LX + 12 + i*(KPI_W+KPI_GAP)
        bx(kx, cy, KPI_W, KPI_H, (255,255,255), (156,163,175), 0.5)
        bx(kx, cy, 3, KPI_H, lclr)
        tx(lbl, kx+8, cy+5,  KPI_W-14, 14, sz=10, bold=True, clr=(75,85,99))
        tx(val, kx+8, cy+22, KPI_W-14, 28, sz=16, bold=True, clr=vclr)
    cy += KPI_H + 12

    # ── 2. Feature Importance Top 5 (L2_fi) / 3. 불량 트렌드 (L3_trend) ─
    show_l2 = "L2_fi" not in hidden_sids
    show_l3 = "L3_trend" not in hidden_sids

    REMAIN_LEFT = BY + SBOX_H - cy - 8
    if show_l2 and show_l3:
        L2_H = (REMAIN_LEFT - 18 - 18) // 2
        L3_H = REMAIN_LEFT - 18 - 18 - L2_H
    elif show_l2:
        L2_H = REMAIN_LEFT - 18; L3_H = 0
    elif show_l3:
        L2_H = 0; L3_H = REMAIN_LEFT - 18
    else:
        L2_H = 0; L3_H = 0

    if show_l2:
        l2_replace = replace_map.get("L2_fi")
        l2_title = l2_replace.get("title") if l2_replace else slabel('L2','Feature Importance Top 5')
        tx(f"{l2_title}", LX+12, cy, LW-24, 16, sz=11, bold=True, clr=(17,24,39))
        cy += 18
        bx(LX+12, cy, LW-24, L2_H, (255,255,255), (156,163,175), 0.5)
        try:
            if l2_replace:
                png = _chart_custom_png(l2_replace, w_px=LW-24, h_px=L2_H)
            else:
                png = _chart_fi_bar_png(features, w_px=LW-24, h_px=L2_H)
            img(png, LX+12, cy, LW-24, L2_H)
        except Exception as _e:
            tx(f"차트 오류: {_e}", LX+16, cy+L2_H//2, LW-32, 18, sz=9, clr=(220,80,80))
        cy += L2_H + 8

    if show_l3:
        l3_replace = replace_map.get("L3_trend")
        l3_title = l3_replace.get("title") if l3_replace else slabel('L3','불량 트렌드')
        tx(f"{l3_title}", LX+12, cy, LW-24, 16, sz=11, bold=True, clr=(17,24,39))
        cy += 18
        bx(LX+12, cy, LW-24, L3_H, (255,255,255), (156,163,175), 0.5)
        try:
            if l3_replace:
                png = _chart_custom_png(l3_replace, w_px=LW-24, h_px=L3_H)
            else:
                png = _chart_trend_png(lot_labels, lot_production, lot_pred_yield,
                                        w_px=LW-24, h_px=L3_H, defect_ppm=lot_defect_ppm,
                                        ww_labels_in=lot_ww_labels)
            img(png, LX+12, cy, LW-24, L3_H)
        except Exception as _e:
            tx(f"차트 오류: {_e}", LX+16, cy+L3_H//2, LW-32, 18, sz=9, clr=(220,80,80))

    # ─── 우칼럼 sbox ───────────────────────────────────────────
    RW = COL_W
    bx(RX, BY, RW, SBOX_H, (255,255,255), (107,114,128), 0.5)
    bx(RX, BY, RW, SHDR_H, (243,244,246), (107,114,128), 0.5)
    bx(RX, BY, 4, SHDR_H, (107,114,128))
    tx(slabel("right_header","[ 불량 예측 현황 ]"), RX+12, BY+8, RW-24, 16, sz=12, bold=True, clr=(17,24,39))

    ry = BY + SHDR_H + 10

    show_r1 = "R1_unit"    not in hidden_sids
    show_r2 = "R2_anomaly" not in hidden_sids
    show_r3 = "R3_scatter" not in hidden_sids

    # ── R1. 대표 Unit (웨이퍼맵 + 정보 테이블, 1:1 분할) ──────
    UNIT_H  = 300 if show_r1 else 0
    INNER_R = RW - 24
    HALF_W  = (INNER_R - 7) // 2   # 1:1 grid with 7px gap
    if show_r1:
        r1_replace = replace_map.get("R1_unit")
        if r1_replace:
            # R1 자리를 custom 차트가 대체
            bx(RX+12, ry, INNER_R, UNIT_H, (255,255,255), (156,163,175), 0.5)
            tx(r1_replace.get("title",""), RX+15, ry+5, INNER_R-8, 14, sz=10, bold=True, clr=(17,24,39))
            try:
                png = _chart_custom_png(r1_replace, w_px=INNER_R-6, h_px=UNIT_H-22)
                img(png, RX+15, ry+22, INNER_R-6, UNIT_H-22)
            except Exception as _e:
                tx(f"차트 오류: {_e}", RX+16, ry+UNIT_H//2, INNER_R-8, 18, sz=9, clr=(220,80,80))
        else:
            # 웨이퍼맵 박스
            bx(RX+12, ry, HALF_W, UNIT_H, (255,255,255), (156,163,175), 0.5)
            tx("불량 위치 웨이퍼맵", RX+15, ry+6, HALF_W-8, 14, sz=10, bold=True, clr=(17,24,39))
            bx(RX+12, ry+22, HALF_W, 1, (209,213,219))
            try:
                # 웨이퍼(정사각형) + 컬러바(세로)를 분리해서 배치 → stretch 없음
                _cbar_w = 26  # 캔버스 폭(라벨 공간 포함), 내부 그라디언트는 ~30%만 차지
                _cbar_gap = 4
                _wmap_side = min(HALF_W - 14 - _cbar_w - _cbar_gap, UNIT_H - 44)
                wmap_png = _chart_wafer_png(wd_dies, wd_x_range, wd_y_range, unit_serial,
                                             w_px=_wmap_side, h_px=_wmap_side)
                cbar_png = _chart_wafer_cbar_png(w_px=_cbar_w, h_px=_wmap_side - 18)
                # 가운데 정렬: 웨이퍼 + cbar 합쳐서 box 중앙
                _total_w = _wmap_side + _cbar_gap + _cbar_w
                _wx = RX + 12 + (HALF_W - _total_w) // 2
                img(wmap_png, _wx, ry+24, _wmap_side, _wmap_side)
                img(cbar_png, _wx + _wmap_side + _cbar_gap, ry+24+9, _cbar_w, _wmap_side - 18)
            except Exception as _e:
                tx(f"웨이퍼맵 오류: {_e}", RX+16, ry+UNIT_H//2, HALF_W-8, 18, sz=9, clr=(220,80,80))
            tx(unit_serial, RX+12, ry+UNIT_H-16, HALF_W, 14, sz=11, bold=True, clr=(17,24,39), align="center")

            # 정보 테이블 (상단 4행 + 포지션 4행)
            INFO_X = RX + 12 + HALF_W + 7
            bx(INFO_X, ry, HALF_W, UNIT_H, (255,255,255), (156,163,175), 0.5)

            pos_health_map = _tu.get("pos_health", {})
            top_rows = [
                ("ufs_serial",  unit_serial,    False),
                ("LOT_ID",      unit_lot,       False),
                ("WAFER_ID",    unit_wafer,     False),
                ("예측 health", _pred_h_disp,   True),
            ]
            POS_HDR_H = 22
            ROW_H = (UNIT_H - POS_HDR_H) // 8
            LBL_W = 105

            for ri, (lbl, val, is_hot) in enumerate(top_rows):
                row_y = ry + ri * ROW_H
                bg = (255,255,255) if ri%2==0 else (248,250,252)
                bx(INFO_X, row_y, HALF_W, ROW_H, bg, (209,213,219), 0.3)
                tx(lbl, INFO_X+6, row_y+4, LBL_W, ROW_H-8, sz=10, bold=True, clr=(55,65,81))
                val_clr = (138,31,31) if is_hot else (17,24,39)
                tx(str(val), INFO_X+LBL_W+6, row_y+4, HALF_W-LBL_W-10, ROW_H-8,
                   sz=11, bold=is_hot, clr=val_clr)

            ph_y = ry + 4*ROW_H
            bx(INFO_X, ph_y, HALF_W, POS_HDR_H, (243,244,246), (209,213,219), 0.3)
            tx("포지션별 기여 (die 평균 대비)", INFO_X+6, ph_y+4, HALF_W-12, POS_HDR_H-8,
               sz=10, bold=True, clr=(17,24,39))

            # 포지션 기여 다이버징 바 — die 평균 중앙선 기준 ±ppm (드릴다운/HTML과 동일)
            _pos_pairs = [(f"P{p}", pos_health_map.get(f"P{p}")) for p in [1, 2, 3, 4]]
            _pnums = [v for _, v in _pos_pairs if isinstance(v, (int, float))]
            if _pnums:
                _pavg = sum(_pnums) / len(_pnums)
                _pmaxabs = max((abs(v - _pavg) for v in _pnums), default=1e-9) or 1e-9
                _ppmax = max(_pnums)
            else:
                _pavg = _pmaxabs = _ppmax = 0
            _pad, _lblw, _ppmw, _devw = 6, 24, 60, 74
            _barx = INFO_X + _pad + _lblw + _ppmw + 4
            _barw = HALF_W - (_pad + _lblw + _ppmw + 4) - _devw - 4
            for ri, (lbl, v) in enumerate(_pos_pairs):
                row_y = ph_y + POS_HDR_H + ri * ROW_H
                bg = (255,255,255) if ri%2==0 else (248,250,252)
                bx(INFO_X, row_y, HALF_W, ROW_H, bg, (209,213,219), 0.3)
                tx(lbl, INFO_X+_pad, row_y+4, _lblw, ROW_H-8, sz=10, bold=True, clr=(55,65,81))
                if not isinstance(v, (int, float)):
                    tx("-", INFO_X+_pad+_lblw, row_y+4, _ppmw, ROW_H-8, sz=10, clr=(156,163,175))
                    continue
                ppm = round(v * 1e6); dev = round((v - _pavg) * 1e6)
                tx(f"{ppm:,} ppm", INFO_X+_pad+_lblw, row_y+4, _ppmw, ROW_H-8,
                   sz=9, bold=True, clr=(17,24,39), align="right")
                bar_yc = row_y + ROW_H // 2
                bx(_barx, bar_yc-5, _barw, 10, (241,245,249), (226,232,240), 0.3)   # 트랙
                _cxb = _barx + _barw // 2
                bx(_cxb-1, bar_yc-7, 2, 14, (148,163,184))                          # 중앙선=die평균
                _sidew = int(abs(v - _pavg) / _pmaxabs * (_barw // 2)) if _pmaxabs else 0
                if (v - _pavg) >= 0:
                    if _sidew > 0: bx(_cxb, bar_yc-3, _sidew, 6, (220,38,38))        # 오른쪽=평균↑(빨강)
                    _dclr = (185,28,28)
                else:
                    if _sidew > 0: bx(_cxb-_sidew, bar_yc-3, _sidew, 6, (148,163,184))  # 왼쪽=평균↓(회색)
                    _dclr = (100,116,139)
                _dtxt = ("+" if dev >= 0 else "") + f"{dev:,}" + (" ◀최다" if v >= _ppmax else "")
                tx(_dtxt, INFO_X+HALF_W-_devw-4, row_y+4, _devw, ROW_H-8,
                   sz=9, bold=True, clr=_dclr, align="right")

        ry += UNIT_H + 10

    # ── R2 Anomaly + R3 분포 비교 그리드 ──────────────────────
    REMAIN_H  = BY + SBOX_H - ry - 4
    PANEL_GAP = 7
    HDR_H_PANEL = 24

    if show_r2 and show_r3:
        R2_W = (RW - 24 - PANEL_GAP) // 2
        R3_X_OFFSET = R2_W + PANEL_GAP
        R3_W = R2_W
    elif show_r2:
        R2_W = RW - 24; R3_X_OFFSET = 0; R3_W = 0
    elif show_r3:
        R2_W = 0; R3_X_OFFSET = 0; R3_W = RW - 24
    else:
        R2_W = 0; R3_X_OFFSET = 0; R3_W = 0

    AX = RX + 12

    # ── R2 SHAP 영향도 패널 (HTML과 동일) ───────────────────────
    if show_r2:
        r2_replace = replace_map.get("R2_anomaly")
        try:
            _n_shap = int(report_data.get("shap_top_n", 5))
            _shap_items = _shap_red_items(report_data, _n_shap)  # 빨강(양수 SHAP=불량↑)만
        except Exception:
            _shap_items = []
        bx(AX, ry, R2_W, REMAIN_H, (255,254,248), (156,163,175), 0.5)
        bx(AX, ry, R2_W, HDR_H_PANEL, (243,244,246), (156,163,175), 0.5)
        r2_title = r2_replace.get("title") if r2_replace else f"SHAP 영향도 Top {len(_shap_items)}"
        tx(r2_title, AX+8, ry+5, R2_W-16, 16, sz=11, bold=True, clr=(17,24,39))
        try:
            if r2_replace:
                png = _chart_custom_png(r2_replace, w_px=R2_W-6, h_px=REMAIN_H-HDR_H_PANEL-4)
            else:
                png = _chart_shap_bar_png(_shap_items, w_px=R2_W-6, h_px=REMAIN_H-HDR_H_PANEL-4)
            img(png, AX+3, ry+HDR_H_PANEL+2, R2_W-6, REMAIN_H-HDR_H_PANEL-4)
        except Exception as _e:
            tx(f"차트 오류: {_e}", AX+6, ry+REMAIN_H//2, R2_W-12, 18, sz=9, clr=(220,80,80))

    # ── R3 피처 정상/불량 분포 비교 패널 ─────────────────────
    if show_r3:
        r3_replace = replace_map.get("R3_scatter")
        # 라벨 = SHAP 빨강(양수) 1등 피처. 분포 그림은 그 피처의 실제 분포를 쓰되,
        # X594는 X727 분포로 치환(하드코딩) — 라벨은 X594 유지.
        _user_fdc = report_data.get("feat_dist_compare") or {}
        _label_feat = _top_red_shap_feature(report_data) or _user_fdc.get("feature", "")
        _DIST_SUB = {"X594": "X727"}
        _data_feat = _DIST_SUB.get(_label_feat, _label_feat)
        try:
            _fdc = _dashboard_feat_dist(feature=_data_feat) or _user_fdc
        except Exception:
            _fdc = _user_fdc
        fdc_feature   = _label_feat or _fdc.get("feature", "")
        fdc_labels    = _fdc.get("labels", [])
        fdc_normal    = _fdc.get("normal", [])
        fdc_danger    = _fdc.get("danger", [])
        fdc_threshold = _fdc.get("threshold")
        if not fdc_feature and features:
            fdc_feature = features[0].get("feature", "")

        FIX = AX + R3_X_OFFSET
        bx(FIX, ry, R3_W, REMAIN_H, (255,255,255), (156,163,175), 0.5)
        bx(FIX, ry, R3_W, HDR_H_PANEL, (243,244,246), (156,163,175), 0.5)
        r3_title = r3_replace.get("title") if r3_replace else f"피처 정상/불량 분포 · {fdc_feature}"
        tx(r3_title, FIX+8, ry+5, R3_W-16, 16, sz=11, bold=True, clr=(17,24,39))

        try:
            if r3_replace:
                png = _chart_custom_png(r3_replace, w_px=R3_W-6, h_px=REMAIN_H-HDR_H_PANEL-4)
            else:
                png = _chart_fdc_line_png(fdc_labels, fdc_normal, fdc_danger, fdc_threshold,
                                           fdc_feature, w_px=R3_W-6, h_px=REMAIN_H-HDR_H_PANEL-4)
            img(png, FIX+3, ry+HDR_H_PANEL+2, R3_W-6, REMAIN_H-HDR_H_PANEL-4)
        except Exception as _e:
            tx(f"차트 오류: {_e}", FIX+6, ry+REMAIN_H//2, R3_W-12, 18, sz=9, clr=(220,80,80))

    # ─── 푸터 ──────────────────────────────────────────────────
    has_extra_page = bool(extra_sections) or bool(commentary_list)
    total_pages = 2 if has_extra_page else 1

    FTR_Y = SH - FTR_H
    bx(0, FTR_Y, SW, FTR_H, (241,245,249))
    bx(0, FTR_Y, SW, 1, (107,114,128))
    tx("We Do Technology | SK hynix", 14, FTR_Y+4, 260, 14,
       sz=10, bold=True, clr=(17,24,39))
    tx(f"{today_str}", SW//2-220, FTR_Y+4, 440, 14,
       sz=10, clr=(75,85,99), align="center")

    # ─── 2번째 슬라이드: 추가 차트 + 메모 (있을 때만) ───────────
    if has_extra_page:
        slide2 = prs.slides.add_slide(prs.slide_layouts[6])

        def bx2(x, y, w, h, fill, line=None, lpt=0.5):
            return _box(slide2, x, y, w, h, fill, line, lpt, dpi=DPI)
        def tx2(text, x, y, w, h, sz=11, bold=False, clr=(32,36,42), align="left", wrap=True):
            return _txt(slide2, text, x, y, w, h, sz, bold, clr, align, wrap, DPI)
        def img2(png_bytes, x, y, w, h):
            return _png_shape(slide2, png_bytes, x, y, w, h, DPI)

        # 헤더
        bx2(0, 0, SW, SH, (255,255,255))
        bx2(0, 0, SW, TOPBAR_H, (255,255,255))
        bx2(0, TOPBAR_H-2, SW, 2, (30,58,138))
        tx2(f"발행일자: {today_str}", 18, 9, 220, 20, sz=12, bold=True, clr=(51,78,118))
        tx2(f"{report_title} (추가)", 0, 7, SW, 24, sz=14, bold=True, clr=(15,23,42), align="center")

        # 본문
        body_y = TOPBAR_H + 16
        cur_y = body_y

        # 메모 영역
        for it in commentary_list:
            note_h = 56
            bx2(20, cur_y, SW-40, note_h, (255,251,235), (253,224,71), 0.75)
            tx2(f"📝 {it.get('title','메모')}", 30, cur_y+6, SW-60, 18,
                sz=12, bold=True, clr=(146,64,14))
            tx2(_strip_html(it.get('content','')), 30, cur_y+26, SW-60, note_h-30,
                sz=11, clr=(30,41,59))
            cur_y += note_h + 8

        # 추가 차트 — 2열 그리드, 각 셀 600×200
        if extra_sections:
            COL_W2 = (SW - 40 - 20) // 2  # 두 열 + 20 gap
            CARD_H2 = 220
            for idx, sec in enumerate(extra_sections):
                col = idx % 2
                row = idx // 2
                cx2 = 20 + col * (COL_W2 + 20)
                cy2 = cur_y + row * (CARD_H2 + 10)
                if cy2 + CARD_H2 > SH - FTR_H - 10:
                    break  # 페이지 초과 시 중단
                bx2(cx2, cy2, COL_W2, CARD_H2, (255,255,255), (156,163,175), 0.5)
                tx2(sec.get("title","섹션"), cx2+8, cy2+6, COL_W2-16, 16,
                    sz=12, bold=True, clr=(17,24,39))
                try:
                    png = _chart_custom_png(sec, w_px=COL_W2-12, h_px=CARD_H2-28)
                    img2(png, cx2+6, cy2+24, COL_W2-12, CARD_H2-28)
                except Exception as _e:
                    tx2(f"차트 오류: {_e}", cx2+10, cy2+CARD_H2//2, COL_W2-20, 18, sz=9, clr=(220,80,80))

        # 푸터
        bx2(0, FTR_Y, SW, FTR_H, (241,245,249))
        bx2(0, FTR_Y, SW, 1, (107,114,128))
        tx2("We Do Technology | SK hynix", 14, FTR_Y+4, 260, 14, sz=10, bold=True, clr=(17,24,39))
        tx2(f"{today_str}", SW//2-220, FTR_Y+4, 440, 14,
            sz=10, clr=(75,85,99), align="center")

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.read()


def markdown_to_pptx(markdown: str) -> bytes:
    """하위 호환용."""
    return build_pptx({
        "meta": {"title": "Field Health 불량 예측 분석 보고서"},
        "scan": {}, "importance": {}, "analysis": {}, "actions": [],
    })


def build_html(report_data: dict) -> str:
    """
    Chart.js 기반 HTML 보고서 (품질불량예측보고서 양식).
    헤더: 발행일자 / 제목 / 대외비 + 전체너비 알림 배너
    왼쪽: 모델링 현황 보고 / 오른쪽: 불량 유닛 분석 현황
    """
    import json as _json
    from tools import get_lot_trend_data

    meta         = report_data.get("meta", {})
    scan         = report_data.get("scan", {})
    importance   = report_data.get("importance", {})
    analysis     = report_data.get("analysis", {})
    features     = importance.get("features", [])
    top_features = analysis.get("top_features", [])
    compare_grp  = analysis.get("compare_group", "grade4")

    chart_params    = report_data.get("chart_params", {})
    shap_top_n      = int(chart_params.get("top_n", 10)) if chart_params.get("chart") == "shap" else 10
    lot_top_n       = int(chart_params.get("top_n", 20)) if chart_params.get("chart") == "lot"  else 20
    table_params    = report_data.get("table_params", {})
    shap_hide_cols  = set(table_params.get("shap_hide_cols", []))
    custom_sections = report_data.get("custom_sections", [])
    commentary_list = report_data.get("commentary", [])

    val_rmse   = meta.get("val_rmse",  "0.005699")
    test_rmse  = meta.get("test_rmse", "0.008427")
    model_nm   = meta.get("model",     "Stacking Ensemble")
    title      = meta.get("title",     "품질불량예측보고서")
    today      = datetime(2026, 6, 11)   # 오늘 고정
    today_str  = today.strftime("%Y. %m. %d")
    scan_total  = scan.get("total_units",   "-")
    scan_high   = scan.get("grade1_count", scan.get("high_count",  "-"))
    scan_ratio  = scan.get("grade1_ratio", scan.get("high_ratio",  "-"))
    top_lot     = scan.get("top_lot",      "-")
    top_wafer   = scan.get("top_wafer",    {})

    # ── 배너: ppm delta 실데이터 (meta로 덮어쓰기 가능)
    ppm_delta   = report_data.get("ppm_delta", {})
    delta_val   = ppm_delta.get("delta", 0)
    delta_str   = f"▲{abs(delta_val):,.0f}" if delta_val >= 0 else f"▼{abs(delta_val):,.0f}"
    delta_color = "#EF4444" if delta_val >= 0 else "#16A34A"
    try:
        _default_alert = ", ".join(it["feature"] for it in _shap_red_items(report_data, 1))   # SHAP 차트 상위 1개(빨강=불량↑)
    except Exception:
        _default_alert = ", ".join(ppm_delta.get("top_features", [f.get("feature","") for f in top_features[:2]])) or "-"
    alert_features = meta.get("alert_features", _default_alert)

    # ── 커스텀 텍스트 (에이전트 수정 가능)
    report_title   = meta.get("report_title",  "Field Health 불량 예측 분석 보고서")
    summary_title  = meta.get("summary_title", f"전 주 대비 품질 불량 &nbsp;<span style=\"color:{delta_color}\">{delta_str} ppm</span>&nbsp; <span style=\"font-size:17px;font-weight:600;color:#555\">{'열화' if delta_val >= 0 else '개선'}</span>")
    summary_sub    = meta.get("summary_sub",   f"원인 WT Parameter&nbsp;<span style=\"background:#fef3c7;color:#92400e;padding:1px 7px;font-size:15px;font-weight:800\">{alert_features}</span>&nbsp;이상 → inline 참원인 도출 요청")

    # ── 섹션 소제목 (에이전트 수정 가능)
    _sl = meta.get("section_labels", {})
    slabel = lambda key, default: _sl.get(key, default)

    # ── L2 주차별 불량 트렌드 (dashboard_units.csv 기반, 대시보드 동일 기준)
    weekly_yield   = report_data.get("weekly_yield_trend", {})
    lot_labels     = weekly_yield.get("labels",     [])
    lot_production = weekly_yield.get("production", [])
    lot_pred_yield = weekly_yield.get("pred_yield", [])
    lot_defect_ppm = weekly_yield.get("defect_ppm", [])
    lot_ww_labels  = weekly_yield.get("ww_labels",  [])   # 'N월 N주차' (대시보드 앵커)
    lot_date_labels = weekly_yield.get("date_labels", [])

    # ── L3: 상위 2개 피처 scatter (임계선 포함)
    feat_scatter = report_data.get("feat_scatter", {})
    fs1 = feat_scatter.get("feat1", {})
    fs2 = feat_scatter.get("feat2", {})
    fs1_name      = fs1.get("name", "Feat1")
    fs2_name      = fs2.get("name", "Feat2")
    fs1_threshold = fs1.get("threshold")
    fs2_threshold = fs2.get("threshold")
    fs1_high = fs1.get("pts_high", [])
    fs1_med  = fs1.get("pts_med",  [])
    fs2_high = fs2.get("pts_high", [])
    fs2_med  = fs2.get("pts_med",  [])

    # ── 피처임포턴스 비율 (분모=전체 X피처 gain 합 — 대시보드 ProcessFactor와 동일)
    total_gain = _fi_total_gain_all() or sum(f.get("lgbm_gain", 0) or 0 for f in features) or 1
    max_gain   = max((f.get("lgbm_gain", 0) or 0 for f in features), default=1) or 1
    fi_ratio_rows = ""
    for i, f in enumerate(features):
        fname = f.get("feature","")
        gain  = f.get("lgbm_gain", 0) or 0
        pct   = round(gain / total_gain * 100, 1)
        bar_w = max(int(gain / max_gain * 100), 3)
        fi_ratio_rows += (
            f'<div class="fi-row">'
            f'<div style="width:18px;font-size:12px;color:#6f756d;text-align:right;font-weight:800;flex-shrink:0">{i+1}</div>'
            f'<div style="width:80px;font-size:13px;font-weight:800;color:#202832;white-space:nowrap;overflow:hidden;text-overflow:ellipsis;font-family:Consolas,monospace;flex-shrink:0" title="{fname}">{fname}</div>'
            f'<div style="flex:1;height:12px;background:#eee9df;border:1px solid #c9c0b1;overflow:hidden;min-width:0">'
            f'<div style="width:{bar_w}%;height:100%;background:#374151"></div>'
            f'</div>'
            f'<div style="width:38px;font-size:12px;color:#374151;text-align:right;font-weight:900;flex-shrink:0">{pct}%</div>'
            f'</div>'
        )

    # ── 어노멀리 피처 (정상/불량 바, grade1 vs grade4 실데이터)
    # (anomaly_stats 계산 제거됨 — R2 패널은 SHAP 영향도로 대체)
    def _short_val(v):
        try:
            n = float(v)
        except (TypeError, ValueError):
            return str(v)
        a = abs(n)
        return f"{n:.0f}" if a >= 1000 else f"{n:.2f}" if a >= 1 else f"{n:.4f}"

    # ── 3a: R2 'SHAP 영향도' — 빨강(양수 SHAP=불량↑)만 표시 ──
    try:
        _n_shap = int(report_data.get("shap_top_n", 5))
        _shap_items = _shap_red_items(report_data, _n_shap)
    except Exception:
        _shap_items = []
    _shap_n = len(_shap_items)
    anomaly_rows = ""
    _smax = max((it["mag"] for it in _shap_items), default=1e-9) or 1e-9
    for it in _shap_items:
        _w = min(100, max(2, it["mag"] / _smax * 100))
        _ppm = it["mag"] * 1e6
        _cls = "danger" if it.get("signed", 0) >= 0 else "normal"
        _dir = "▲ 불량↑" if it.get("signed", 0) >= 0 else "▼ 불량↓"
        _dclr = "#b91c1c" if it.get("signed", 0) >= 0 else "#1d4ed8"
        anomaly_rows += (
            f'<div class="anom-card">'
            f'<div class="anom-top">'
            f'<div class="anom-name" title="{it["feature"]}">{it["feature"]}</div>'
            f'<span style="color:{_dclr};font-size:11px;font-weight:900">{_dir} {_ppm:,.0f}ppm</span>'
            f'</div>'
            f'<div class="anom-row">'
            f'<div class="anom-track"><div class="anom-fill {_cls}" style="width:{_w:.1f}%"></div></div>'
            f'</div>'
            f'</div>'
        )

    # ── 더미 섹션 감지 (실데이터 없으면 더미로 표시)
    _is_dummy_l2   = not report_data.get("weekly_yield_trend", {}).get("labels")
    _is_dummy_l3   = not report_data.get("lot_defect_counts") and not report_data.get("pred_ppm_trend", {}).get("labels")
    _is_dummy_l4   = not report_data.get("feat_vs_health", {}).get("high_pts") and not report_data.get("shap_trend", {}).get("high_pts")
    _is_dummy_r1   = not report_data.get("top_unit") and not report_data.get("wafer_die", {}).get("dies")
    _is_dummy_r1b  = not report_data.get("top_unit", {}).get("pos_health") and not report_data.get("top_unit", {}).get("pos_feat_vals")
    def _dummy_attr(flag): return ' data-dummy="1"' if flag else ''

    # ── 포지션별 WT 이상 비율 실데이터
    pos_defect     = report_data.get("pos_defect", {})
    pos_labels     = pos_defect.get("labels",     ["P1","P2","P3","P4"])
    pos_high_ratio = pos_defect.get("high_ratio", [0, 0, 0, 0])
    pos_med_ratio  = pos_defect.get("med_ratio",  [0, 0, 0, 0])
    pos_low_ratio  = pos_defect.get("low_ratio",  [0, 0, 0, 0])
    pos_feat1      = pos_defect.get("feat1", features[0].get("feature","F1") if features else "F1")
    pos_feat2      = pos_defect.get("feat2", features[1].get("feature","F2") if len(features)>1 else "F2")

    # 포지션 어노멀리 스타일 HTML (각 포지션 × feat1이상/feat2이상/정상 막대)
    pos_rows = ""
    for i, lbl in enumerate(pos_labels):
        h = pos_high_ratio[i] if i < len(pos_high_ratio) else 0
        m = pos_med_ratio[i]  if i < len(pos_med_ratio)  else 0
        l = pos_low_ratio[i]  if i < len(pos_low_ratio)  else 0
        pos_rows += (
            f'<tr>'
            f'<td style="font-size:13px;width:28px;padding:3px 5px;vertical-align:middle">{lbl}</td>'
            f'<td style="padding:3px 5px">'
            f'<div style="margin-bottom:2px">'
            f'<div style="font-size:11px;color:#EF4444;margin-bottom:1px">{pos_feat1} 이상 {h}%</div>'
            f'<div style="height:6px;border-radius:2px;background:#EF4444;width:{min(int(h*4),100)}%"></div>'
            f'</div>'
            f'<div style="margin-bottom:2px">'
            f'<div style="font-size:11px;color:#F59E0B;margin-bottom:1px">{pos_feat2} 이상 {m}%</div>'
            f'<div style="height:6px;border-radius:2px;background:#F59E0B;width:{min(int(m*4),100)}%"></div>'
            f'</div>'
            f'<div>'
            f'<div style="font-size:11px;color:#16A34A;margin-bottom:1px">정상 {l}%</div>'
            f'<div style="height:6px;border-radius:2px;background:#16A34A;width:{min(int(l*4),100)}%"></div>'
            f'</div>'
            f'</td>'
            f'</tr>'
        )

    # ── 대표 Unit (reg_pred 최고 unit 실데이터)
    _tu = report_data.get("top_unit", {})
    _run_id   = _tu.get("run_id", "-")
    _wafer_no = _tu.get("wafer_no", "-")
    # 생산 일자: 오늘로부터 run_id 기반 역산 (val LOT: 약 1주 전~5주 전)
    _run_id_int = int(_run_id) if isinstance(_run_id, (int,float)) and str(_run_id) != "-" else 53
    _val_lots_sorted = sorted(range(29, 57))  # val: run_id 29~56
    _lot_idx = _val_lots_sorted.index(_run_id_int) if _run_id_int in _val_lots_sorted else 0
    # 생산일자·예측일자 모두 오늘(6/11)로 고정
    _prod_date = today_str
    _insp_date = today_str
    _pred_health_raw = _tu.get("pred_health", _tu.get("pred_ppm", 0))
    try:
        _pred_health_f = float(_pred_health_raw)
    except (TypeError, ValueError):
        _pred_health_f = 0.0
    # pred_health가 ppm 단위로 저장된 경우 → health 단위로 변환 후 ppm 표시
    _pred_health_frac = (_pred_health_f / 1e6) if _pred_health_f > 1 else _pred_health_f
    _pred_health_display = f"{round(_pred_health_frac * 1e6):,} ppm"
    dummy_unit = {
        "serial":        _tu.get("serial", "S38369"),
        "pred_health":   _pred_health_display,
        "pred_ppm":      f'{_tu.get("pred_ppm", 5241.9):,.1f}' if isinstance(_tu.get("pred_ppm"), (int,float)) else str(_tu.get("pred_ppm", "5241.9")),
        "lot":           _tu.get("run_id", _run_id),
        "wafer":         _tu.get("wafer_no", _wafer_no),
        "risk":          _tu.get("risk", "HIGH"),
        "prod_date":     _prod_date,
        "insp_date":     today_str,
    }
    # 포지션별 기여 — die 평균 기준 다이버징 바 (드릴다운과 동일 표현)
    # pos_health[Pn] = wafer_map die-level pred(= 유닛 예측에 대한 die 기여). ppm 변환 후
    # 4 die 평균을 중앙선으로, 평균보다 많이 기여=오른쪽(빨강)/적게=왼쪽(회색).
    pos_health_map = _tu.get("pos_health", {})
    pos_feat_vals  = _tu.get("pos_feat_vals", {})
    pos_health_rows = ""
    _pos_pairs = [(p, pos_health_map.get(p)) for p in ["P1", "P2", "P3", "P4"]]
    _pos_nums  = [v for _, v in _pos_pairs if isinstance(v, (int, float))]
    if _pos_nums:
        _avg    = sum(_pos_nums) / len(_pos_nums)
        _maxabs = max((abs(v - _avg) for v in _pos_nums), default=1e-9) or 1e-9
        _pmax   = max(_pos_nums)
        for p, v in _pos_pairs:
            if not isinstance(v, (int, float)):
                pos_health_rows += (
                    f'<div style="display:flex;align-items:center;gap:4px;height:26px;overflow:hidden">'
                    f'<div class="unit-lbl" style="min-width:18px;flex-shrink:0">{p}</div>'
                    f'<div style="color:#9ca3af;flex:1">-</div></div>'
                )
                continue
            ppm  = round(v * 1e6)
            dev  = round((v - _avg) * 1e6)
            side = round(abs(v - _avg) / _maxabs * 50)
            ismax = v >= _pmax
            if (v - _avg) >= 0:
                fill = (f'<div style="position:absolute;left:50%;top:2px;bottom:2px;width:{side}%;'
                        f'background:linear-gradient(90deg,#FCA5A5,#DC2626);border-radius:0 2px 2px 0"></div>')
                devcolor = "#B91C1C"
            else:
                fill = (f'<div style="position:absolute;right:50%;top:2px;bottom:2px;width:{side}%;'
                        f'background:linear-gradient(90deg,#94A3B8,#CBD5E1);border-radius:2px 0 0 2px"></div>')
                devcolor = "#64748B"
            track = (
                f'<div style="position:relative;flex:0 0 110px;height:13px;'
                f'background:linear-gradient(90deg,#F1F5F9,#F8FAFC 50%,#F1F5F9);border-radius:3px">'
                f'<div style="position:absolute;left:50%;top:-2px;bottom:-2px;width:2px;background:#94A3B8;'
                f'transform:translateX(-50%);border-radius:1px"></div>{fill}</div>'
            )
            maxtag = '<span style="color:#B91C1C;font-weight:800;font-size:7px;margin-left:2px">◀최다</span>' if ismax else ''
            pos_health_rows += (
                f'<div style="display:flex;align-items:center;gap:4px;height:26px;overflow:hidden">'
                f'<div class="unit-lbl" style="min-width:18px;flex-shrink:0;font-size:11px">{p}</div>'
                f'<div style="font-family:Consolas,monospace;font-size:10px;font-weight:700;color:#111827;'
                f'min-width:40px;text-align:right;flex-shrink:0">{ppm:,}<span style="font-size:7px;color:#6b7280"> ppm</span></div>'
                f'{track}'
                f'<div style="font-family:Consolas,monospace;font-size:9px;font-weight:700;color:{devcolor};'
                f'min-width:36px;text-align:right;flex-shrink:0;white-space:nowrap">{"+" if dev>=0 else ""}{dev:,}{maxtag}</div>'
                f'</div>'
            )
    else:
        # fallback: pos_health 없으면 feature 값 표시 (기존)
        _pf_feats = []
        if pos_feat_vals:
            sample_p = next(iter(pos_feat_vals.values()), {})
            _pf_feats = list(sample_p.keys())[:2]
        if not _pf_feats and features:
            _pf_feats = [features[0].get("feature", "")]
        for p in ["P1", "P2", "P3", "P4"]:
            pdata = pos_feat_vals.get(p, {})
            if pdata and _pf_feats:
                val_display = "  ".join(
                    f'<span style="color:#4b5563;font-size:10px">{f}=</span>'
                    f'<span style="font-family:Consolas,monospace;font-size:11px;font-weight:900;color:#1e3a8a">{pdata.get(f,"?")}</span>'
                    for f in _pf_feats if f in pdata
                )
            else:
                val_display = '<span style="color:#9ca3af">-</span>'
            pos_health_rows += (
                f'<div class="unit-row" style="grid-template-columns:104px 1fr;height:32px;align-items:center">'
                f'<div class="unit-lbl">{p}</div>'
                f'<div class="unit-val">{val_display}</div>'
                f'</div>'
            )

    # ── Lot별 불량 개수 (L3, 1팀 스타일)
    lot_defect = report_data.get("lot_defect_counts", [])  # [{"lot": "Lot 43", "count": 5, "rate": 2.1}, ...]
    if not lot_defect:
        # pred_ppm_trend 에서 변환
        pred_ppm_trend = report_data.get("pred_ppm_trend", {})
        _lbals = pred_ppm_trend.get("labels", [f"Lot {40+i}" for i in range(14)])
        _lhigh = pred_ppm_trend.get("defect_count", pred_ppm_trend.get("high_ppm", []))
        _lrate = pred_ppm_trend.get("defect_rate", [])
        _lhigh_safe = (_lhigh or []) + [0] * max(0, len(_lbals) - len(_lhigh or []))
        _lrate_safe = (_lrate or []) + [0] * max(0, len(_lbals) - len(_lrate or []))
        lot_defect = [{"lot": l, "count": c, "rate": r}
                      for l, c, r in zip(_lbals, _lhigh_safe, _lrate_safe)]

    j_lot_defect_labels = _json.dumps([d["lot"] for d in lot_defect], ensure_ascii=False)
    j_lot_defect_counts = _json.dumps([d.get("count", 0) for d in lot_defect])
    j_lot_defect_rates  = _json.dumps([d.get("rate", 0) for d in lot_defect])

    # ── L4: 피처값 vs 예측 health scatter (feat_vs_health 실데이터)
    fvh         = report_data.get("feat_vs_health", {})
    # shap_trend fallback도 지원
    _shap       = report_data.get("shap_trend", {})
    if fvh.get("high_pts") or fvh.get("normal_pts"):
        shap_feat_name = fvh.get("feature", features[0].get("feature","") if features else "")
        j_shap_high    = _json.dumps(fvh.get("high_pts",  []))
        j_shap_normal  = _json.dumps(fvh.get("normal_pts",[]))
        j_shap_x_label = fvh.get("x_label", shap_feat_name)
        j_shap_y_label = "reg_pred (health)"
        _is_fvh_real   = True
    else:
        shap_feat_name = _shap.get("feature", features[0].get("feature","") if features else "")
        j_shap_high    = _json.dumps(_shap.get("high_pts",  []))
        j_shap_normal  = _json.dumps(_shap.get("normal_pts",[]))
        j_shap_x_label = shap_feat_name
        j_shap_y_label = "SHAP value"
        _is_fvh_real   = False
    j_shap_trend_y  = _json.dumps(_shap.get("mean_trend",[]))
    j_shap_trend_x  = _json.dumps(_shap.get("trend_x",  []))

    # ── R3: LOT별 예측 ppm 트렌드 (HIGH/MED 그룹 평균) — 하위호환
    pred_ppm_trend = report_data.get("pred_ppm_trend", {})
    if pred_ppm_trend.get("labels"):
        r3_labels = pred_ppm_trend["labels"]
        r3_high   = pred_ppm_trend.get("high_ppm", [])
        r3_med    = pred_ppm_trend.get("med_ppm",  [])
    else:
        r3_labels = [f"Lot {40+i}" for i in range(14)]
        r3_high   = []; r3_med = []

    # ── custom_sections 렌더링 (다른 차트와 동일한 박스 스타일)
    def _render_custom_section(sec, idx):
        sec_title = sec.get("title",""); ctype = sec.get("chart_type","bar")
        height = sec.get("height", 120)
        if ctype == "table":
            cols = sec.get("columns",[]); rows = sec.get("rows",[])
            ths = "".join(f"<th>{c}</th>" for c in cols)
            trs = ""
            for ri, row in enumerate(rows):
                cells = row if isinstance(row, list) else [row.get(c,"") for c in cols]
                bg = "#fff" if ri%2==0 else "#F8FAFC"
                trs += "<tr>" + "".join(f'<td style="background:{bg}">{v}</td>' for v in cells) + "</tr>"
            content = f'<div style="max-height:{height}px;overflow:auto;padding:4px"><table style="width:100%;border-collapse:collapse;font-size:11px"><thead><tr style="background:#f3f4f6">{ths}</tr></thead><tbody>{trs}</tbody></table></div>'
        else:
            content = f'<div style="position:relative;height:{height}px;padding:4px"><canvas id="cs_{idx}_chart" style="position:absolute;top:4px;left:4px;right:4px;bottom:4px;width:calc(100% - 8px);height:calc(100% - 8px)"></canvas></div>'
        return (
            f'<div class="cbox ia-target" data-sid="cs_{idx}" data-section="{sec_title}" '
            f'style="margin-bottom:8px;position:relative;border:1px solid #9ca3af;background:#fff;display:flex;flex-direction:column">'
            f'<div style="padding:5px 9px;background:#f3f4f6;border-bottom:1px solid #d1d5db;font-size:12px;font-weight:900;color:#111827">{sec_title}</div>'
            f'{content}</div>'
        )

    def _render_custom_chart_js(sec, idx):
        ctype = sec.get("chart_type", "bar")
        if ctype == "table": return ""
        labels   = _json.dumps(sec.get("labels", []), ensure_ascii=False)
        stacked  = sec.get("stacked", False)
        ds_js = []
        for ds in sec.get("datasets", []):
            color = ds.get("color", "#3B82F6")
            bg    = _json.dumps(ds["colors"]) if "colors" in ds else f'"{color}"'
            ds_js.append(
                f'{{"label":{_json.dumps(ds.get("label",""))},"data":{_json.dumps(ds.get("data",[]))},'
                f'"backgroundColor":{bg},"borderColor":"{color}","borderWidth":2,"tension":0.3,"fill":false}}'
            )
        horizontal = _json.dumps(sec.get("horizontal", False))
        scales_js  = ',"scales":{"x":{"stacked":true},"y":{"stacked":true}}' if stacked else ''
        return (
            f"(function(){{var ctx=document.getElementById('cs_{idx}_chart');if(!ctx)return;"
            f"new Chart(ctx,{{type:{_json.dumps(ctype)},data:{{labels:{labels},"
            f"datasets:[{','.join(ds_js)}]}},"
            f"options:{{indexAxis:{horizontal}?'y':'x',responsive:true,maintainAspectRatio:false,"
            f"plugins:{{legend:{{display:true,position:'top',labels:{{boxWidth:10,font:{{size:14}}}}}}}}"
            f"{scales_js}}}}});}})();"
        )

    def _render_replace_section_js(sec, idx):
        """replace_sid 있는 섹션: 원본 위치에 인라인 주입 (JS DOM 조작)."""
        replace_sid = sec.get("replace_sid", "")
        sec_title   = sec.get("title", "")
        ctype       = sec.get("chart_type", "bar")
        height      = sec.get("height", 120)
        canvas_id   = f"cs_{idx}_chart"
        replace_attr = f'data-replace-for="{replace_sid}"'
        new_sid      = f"cs_{idx}"
        section_html = (
            f'<div class="cbox ia-target" {replace_attr} data-sid="{new_sid}" data-origin-sid="{replace_sid}" data-section="{sec_title}" '
            f'style="margin-bottom:8px;position:relative;border:1px solid #9ca3af;background:#fff;display:flex;flex-direction:column">'
            f'<div style="padding:5px 9px;background:#f3f4f6;border-bottom:1px solid #d1d5db;font-size:12px;font-weight:900;color:#111827">{sec_title}</div>'
            f'<div style="position:relative;height:{height}px;padding:4px">'
            f'<canvas id="{canvas_id}" style="position:absolute;top:4px;left:4px;right:4px;bottom:4px;width:calc(100% - 8px);height:calc(100% - 8px)"></canvas>'
            f'</div></div>'
        )
        html_js = _json.dumps(section_html)
        labels    = _json.dumps(sec.get("labels", []), ensure_ascii=False)
        stacked   = sec.get("stacked", False)
        ds_js = []
        for ds in sec.get("datasets", []):
            color = ds.get("color", "#3B82F6")
            bg    = _json.dumps(ds["colors"]) if "colors" in ds else f'"{color}"'
            ds_js.append(
                f'{{"label":{_json.dumps(ds.get("label",""))},"data":{_json.dumps(ds.get("data",[]))},'
                f'"backgroundColor":{bg},"borderColor":"{color}","borderWidth":2,"tension":0.3,"fill":false}}'
            )
        horizontal = _json.dumps(sec.get("horizontal", False))
        scales_js  = ',"scales":{"x":{"stacked":true},"y":{"stacked":true}}' if stacked else ''
        if ctype == "table":
            chart_init = ""
        else:
            chart_init = (
                'var ctx=document.getElementById(' + _json.dumps(canvas_id) + ');'
                'if(ctx)new Chart(ctx,{type:' + _json.dumps(ctype) + ','
                'data:{labels:' + labels + ',datasets:[' + ','.join(ds_js) + ']},'
                'options:{indexAxis:' + horizontal + '?"y":"x",'
                'responsive:true,maintainAspectRatio:false,'
                'plugins:{legend:{display:true,position:"top",'
                'labels:{boxWidth:10,font:{size:9}}}}' + scales_js + '}});'
            )
        return (
            f'(function(){{'
            f'var prev=document.querySelector(\'[data-replace-for="{replace_sid}"]\');'
            f'if(prev)prev.remove();'
            f'var orig=document.querySelector(\'[data-sid="{replace_sid}"]\');'
            f'if(!orig)return;'
            f'orig.style.display="none";'
            f'var lbl=orig.previousElementSibling;'
            f'if(lbl&&lbl.classList&&lbl.classList.contains("inum"))lbl.style.display="none";'
            f'var wrapper=document.createElement("div");'
            f'wrapper.innerHTML={html_js};'
            f'orig.parentNode.insertBefore(wrapper.firstChild,orig.nextSibling);'
            f'{chart_init}'
            f'}})();'
        )

    # replace_sid 있으면 인라인 주입, 없으면 기존 {left_extra}/{right_extra} 슬롯
    _add_secs     = [(i, s) for i, s in enumerate(custom_sections) if not s.get("replace_sid")]
    _replace_secs = [(i, s) for i, s in enumerate(custom_sections) if s.get("replace_sid")]
    left_extra       = "".join(_render_custom_section(s, i) for i, s in _add_secs if s.get("position") == "left_col")
    right_extra      = "".join(_render_custom_section(s, i) for i, s in _add_secs if s.get("position") == "right_col")
    custom_chart_js  = "\n".join(_render_custom_chart_js(s, i) for i, s in _add_secs)
    replace_chart_js = "\n".join(_render_replace_section_js(s, i) for i, s in _replace_secs)
    commentary_html = "\n".join(
        f'<div style="background:#FFFBEB;border:1px solid #FDE68A;border-radius:8px;padding:10px 14px;margin-bottom:8px">'
        f'<div style="font-size:13px;font-weight:700;color:#92400E;margin-bottom:4px">📝 {it.get("title","")}</div>'
        f'<div style="font-size:13px;color:#1E293B">{it.get("content","")}</div></div>'
        for it in commentary_list
    )

    # ── 웨이퍼맵 die 좌표
    wafer_die   = report_data.get("wafer_die", {})
    wd_dies     = wafer_die.get("dies", [])
    wd_x_range  = wafer_die.get("x_range", [12, 66])
    wd_y_range  = wafer_die.get("y_range", [11, 32])
    wd_serial   = wafer_die.get("serial", dummy_unit["serial"])

    # ── 웨이퍼맵 HTML (pred 연속 그라디언트 색상, DrilldownV2 동일)
    wmap_w, wmap_h = 240, 240
    cx, cy_c = wmap_w / 2, wmap_h / 2
    pad = 14
    r_svg = min(cx, cy_c) - pad

    # die-level pred 기반 색상 스케일 — DrilldownV2와 동일하게 wafer_scale.json 글로벌 기준 사용
    _preds = [d.get("pred") if d.get("pred") is not None else (d.get("pred_ppm", 0) or 0) / 1e6
              for d in wd_dies]
    # wafer_scale.json: 전체 wafer_map.csv 기반 글로벌 스케일 (DrilldownV2의 globalScale과 동일)
    try:
        import json as _json_mod, os as _os_mod
        _scale_path = _os_mod.path.join(_os_mod.path.dirname(__file__), "..", "data", "processed", "wafer_scale.json")
        with open(_os_mod.path.normpath(_scale_path)) as _sf:
            _ws = _json_mod.load(_sf)
        _pred_min  = float(_ws.get("pred_min", 0))
        _pred_max  = float(_ws.get("pred_max", 0.006))
        _threshold = float(_ws.get("threshold", 0.003))
    except Exception:
        # fallback: 현재 wafer 기준
        _pred_min = min(_preds) if _preds else 0
        _pred_max = max(_preds) if _preds else 1
        _sorted_p = sorted(_preds)
        _threshold = _sorted_p[int(len(_sorted_p) * 0.75)] if _sorted_p else (_pred_max * 0.5)

    # DrilldownV2와 동일한 글로벌 좌표 기준 사용 (웨이퍼마다 격자 고정)
    _GLOBAL_X_MIN, _GLOBAL_X_MAX = 12, 66   # GLOBAL_DIE_X_MIN/MAX
    _GLOBAL_Y_MIN, _GLOBAL_Y_MAX = 11, 32   # GLOBAL_DIE_Y_MIN/MAX
    _ref_x_range = _GLOBAL_X_MAX - _GLOBAL_X_MIN + 1  # 55
    _ref_y_range = _GLOBAL_Y_MAX - _GLOBAL_Y_MIN + 1  # 22
    _mx = (_GLOBAL_X_MIN + _GLOBAL_X_MAX) / 2         # centerX = 39
    _my = (_GLOBAL_Y_MIN + _GLOBAL_Y_MAX) / 2         # centerY = 21.5

    # DrilldownV2: cellW = (D / refXRange) * SCALE, D=800, SCALE=0.9
    # 보고서: r_svg 기준으로 동일 비율 적용
    _SCALE = 0.9
    _sx = (r_svg * 2 / _ref_x_range) * _SCALE
    _sy = (r_svg * 2 / _ref_y_range) * _SCALE

    def _wsvgx(x): return cx + (x - _mx) * _sx
    def _wsvgy(y): return cy_c + (y - _my) * _sy

    # die 크기: DrilldownV2와 동일 (cellW × cellH)
    _die_w = _sx
    _die_h = _sy

    die_svgs = ""
    for die in wd_dies:
        px = _wsvgx(die["x"]); py = _wsvgy(die["y"])
        die_pred = die.get("pred") if die.get("pred") is not None else (die.get("pred_ppm", 0) or 0) / 1e6
        is_target = die.get("is_target", False)
        clr = _pred_color(die_pred, _pred_min, _pred_max, _threshold)
        stroke = "#7f1d1d" if is_target else "none"
        sw = "1.5" if is_target else "0"
        die_svgs += (
            f'<rect x="{px-_die_w/2:.1f}" y="{py-_die_h/2:.1f}" '
            f'width="{_die_w:.1f}" height="{_die_h:.1f}" '
            f'fill="{clr}" stroke="{stroke}" stroke-width="{sw}" rx="1"/>'
        )
    wafer_svg = (
        f'<svg width="{wmap_w}" height="{wmap_h}" style="display:block">'
        f'<defs><clipPath id="wafer-clip"><circle cx="{cx:.1f}" cy="{cy_c:.1f}" r="{r_svg:.1f}"/></clipPath></defs>'
        f'<circle cx="{cx:.1f}" cy="{cy_c:.1f}" r="{r_svg:.1f}" fill="#F8FAFC" stroke="#94A3B8" stroke-width="1.5"/>'
        f'<g clip-path="url(#wafer-clip)">{die_svgs}</g>'
        f'<rect x="{cx-3:.1f}" y="{cy_c+r_svg-3:.1f}" width="6" height="6" fill="#CBD5E1" rx="1"/>'
        f'</svg>'
    )

    # ── JSON 직렬화
    j_lot_labels     = _json.dumps(lot_labels, ensure_ascii=False)
    j_lot_production = _json.dumps(lot_production)
    j_lot_pred_yield = _json.dumps(lot_pred_yield)
    j_lot_defect_ppm = _json.dumps(lot_defect_ppm)
    j_lot_ww_labels  = _json.dumps(lot_ww_labels, ensure_ascii=False)
    j_lot_date_labels = _json.dumps(lot_date_labels, ensure_ascii=False)
    j_fs1_high       = _json.dumps(fs1_high)
    j_fs1_med        = _json.dumps(fs1_med)
    j_fs1_threshold  = _json.dumps(fs1_threshold)
    j_fs2_high       = _json.dumps(fs2_high)
    j_fs2_med        = _json.dumps(fs2_med)
    j_fs2_threshold  = _json.dumps(fs2_threshold)
    j_r3_labels      = _json.dumps(r3_labels, ensure_ascii=False)
    j_r3_high        = _json.dumps(r3_high)
    j_r3_med         = _json.dumps(r3_med)
    # Feature Importance — 대시보드 기준(X피처, gain 순) Top 5
    _fi_top4 = features[:5]
    j_fi_top4_labels = _json.dumps([f.get("feature","") for f in _fi_top4], ensure_ascii=False)
    _fi_total = _fi_total_gain_all() or sum(f.get("lgbm_gain", 0) or 0 for f in features) or 1
    j_fi_top4_values = _json.dumps([round((f.get("lgbm_gain", 0) or 0) / _fi_total * 100, 2) for f in _fi_top4])
    # 이상 피처 분포 scatter 용 (fs1_high/fs1_med 재사용, 원본 x/y 그대로)
    _j_feat_scatter_high = _json.dumps(fs1_high)
    _j_feat_scatter_med  = _json.dumps(fs1_med)

    # R3: 피처 정상/위험 분포 — 라벨 = SHAP 빨강(양수) 1등 피처.
    #   분포 그림은 그 피처의 실제 분포를 쓰되, X594는 X727 분포로 치환(하드코딩). 라벨은 X594 유지.
    _user_fdc = report_data.get("feat_dist_compare") or {}
    _label_feat = _top_red_shap_feature(report_data) or _user_fdc.get("feature", "")
    _DIST_SUB = {"X594": "X727"}
    _data_feat = _DIST_SUB.get(_label_feat, _label_feat)
    try:
        fdc = _dashboard_feat_dist(feature=_data_feat) or _user_fdc
    except Exception:
        fdc = _user_fdc
    fdc = dict(fdc); fdc["feature"] = _label_feat or fdc.get("feature", "")   # 라벨만 빨강1등
    j_fdc_feature   = _json.dumps(fdc.get("feature", ""))
    j_fdc_labels    = _json.dumps(fdc.get("labels", []))
    j_fdc_normal    = _json.dumps(fdc.get("normal", []))
    j_fdc_danger    = _json.dumps(fdc.get("danger", []))
    j_fdc_threshold = _json.dumps(fdc.get("threshold"))

    # PPM 계산 (KPI용) — 전체 unit 평균 reg_pred를 ppm으로 환산 (grade1 비율 아님)
    try:
        from tools import get_mean_pred_ppm
        _ppm_val = get_mean_pred_ppm()
        _ppm_str = f"{_ppm_val:,}" if _ppm_val > 0 else "-"
    except Exception:
        _ppm_str = "-"

    html = f"""<!DOCTYPE html>
<html lang="ko">
<head>
<meta charset="UTF-8">
<title>{title}</title>
<script src="https://cdn.jsdelivr.net/npm/chart.js@4.4.2/dist/chart.umd.min.js"></script>
<style>
*{{box-sizing:border-box;margin:0;padding:0}}
html,body{{width:1280px;height:720px;overflow:hidden;background:#e7e7e7}}
body{{font-family:'Malgun Gothic','Segoe UI',Arial,sans-serif;font-weight:600;color:#20242a;font-size:13px}}
.slide{{width:1280px;height:720px;background:#fff;border:1px solid #6b7280;box-shadow:0 10px 24px rgba(31,41,55,.12);display:flex;flex-direction:column;overflow:hidden}}
.s-topbar{{display:grid;grid-template-columns:200px 1fr 200px;align-items:center;height:36px;padding:0 20px;border-bottom:2px solid #1e3a8a;flex-shrink:0}}
.s-issue{{font-size:13px;font-weight:700;color:#334e76}}
.s-caption{{font-size:19px;font-weight:900;color:#0f172a;text-align:center}}
.conf{{background:#dc2626;color:#fff;font-size:13px;font-weight:800;padding:4px 13px;border:1px solid #b91c1c}}
.s-conf-wrap{{display:flex;justify-content:flex-end}}
.s-summary{{background:#fff7ed;border-bottom:1px solid #fed7aa;padding:4px 16px;text-align:center;height:44px;box-sizing:border-box;flex-shrink:0}}
.s-title{{font-size:18px;font-weight:900;color:#111827;line-height:1.2}}
.s-subtitle{{font-size:13px;font-weight:700;color:#1a3a5c;margin-top:1px}}
.s-body{{padding:6px 10px 4px;height:614px;box-sizing:border-box;overflow:hidden;flex-shrink:0}}
.cols{{display:grid;grid-template-columns:1fr 1fr;gap:10px;height:100%}}
.sbox{{border:1px solid #6b7280;background:#fff;display:flex;flex-direction:column;overflow:hidden;height:604px}}
.shdr{{background:#f3f4f6;color:#111827;border-bottom:1px solid #6b7280;padding:4px 11px;font-size:13px;font-weight:900;position:relative;height:28px;box-sizing:border-box;flex-shrink:0}}
.shdr::before{{content:'';position:absolute;left:0;top:0;bottom:0;width:4px;background:#6b7280}}
.sbdy{{padding:5px 9px;height:576px;box-sizing:border-box;overflow-y:auto;overflow-x:hidden}}
.sbdy.left-sbdy{{display:flex;flex-direction:column}}
.sbdy.right-sbdy{{display:flex;flex-direction:column}}
.kpi-cards{{display:grid;grid-template-columns:1fr 1fr 1fr;gap:5px;margin-bottom:5px;height:58px;box-sizing:border-box;flex-shrink:0}}
.kpi-card{{background:#fff;border:1px solid #9ca3af;padding:5px 8px 4px 12px;position:relative;overflow:hidden}}
.kpi-card::before{{content:'';position:absolute;left:0;top:0;bottom:0;width:3px}}
.kpi-card.navy::before{{background:#1e3a8a}}.kpi-card.dark::before{{background:#374151}}.kpi-card.amber::before{{background:#b45309}}
.kpi-lbl{{font-size:12px;color:#4b5563;font-weight:800;margin-bottom:1px}}
.kpi-val{{font-size:18px;font-weight:900;line-height:1.1}}
.kpi-val.navy{{color:#1e3a8a}}.kpi-val.dark{{color:#374151}}.kpi-val.amber{{color:#b45309}}
.kpi-sub{{font-size:12px;color:#4b5563;font-weight:700;margin-top:1px}}
.inum{{font-size:13px;font-weight:900;color:#111827;margin:4px 0 2px;height:18px;box-sizing:border-box;flex-shrink:0}}
.cbox{{border:1px solid #9ca3af;background:#fff;margin-bottom:5px;flex-shrink:0}}
.cbox-body{{padding:3px 5px;position:relative}}
.unit-main{{display:grid;grid-template-columns:1fr 1fr;gap:7px;margin-bottom:5px;flex-shrink:0;align-items:stretch}}
.wafer-box{{border:1px solid #9ca3af;background:#fff;display:flex;flex-direction:column;align-items:center;padding:5px;gap:2px;height:100%}}
.wafer-box-title{{font-size:12px;font-weight:900;color:#111827;align-self:stretch;border-bottom:1px solid #d1d5db;padding-bottom:3px;margin-bottom:1px}}
.unit-tbl{{border:1px solid #9ca3af;overflow:hidden;background:#fff}}
.unit-row{{display:grid;grid-template-columns:72px 1fr;border-bottom:1px solid #d1d5db}}
.unit-row:last-child{{border-bottom:none}}
.unit-row:nth-child(even){{background:#f8fafc}}
.unit-lbl{{padding:4px 6px;font-size:13px;font-weight:900;color:#111827}}
.unit-val{{padding:4px 6px;font-size:12px;font-weight:900;font-family:Consolas,monospace;color:#111827}}
.unit-val.hot{{color:#8a1f1f;font-size:13px}}
.anom-panel{{border:1px solid #9ca3af;background:#fffef8;display:flex;flex-direction:column;overflow:hidden}}
.anom-hdr{{background:#f3f4f6;border-bottom:1px solid #9ca3af;padding:4px 8px;font-size:13px;font-weight:900;flex-shrink:0}}
.anom-list{{padding:2px;display:flex;flex-direction:column;gap:1px;flex:1;overflow:hidden}}
.anom-card{{border:1px solid #d1d5db;background:#fff;padding:2px 6px;flex:1;min-height:0;display:flex;flex-direction:column;overflow:hidden}}
.anom-card:nth-child(even){{background:#f8fafc}}
.anom-top{{display:flex;align-items:center;justify-content:space-between;gap:6px;margin-bottom:1px;flex-shrink:0}}
.anom-name{{color:#111827;font:900 11px/1.2 Consolas,monospace;white-space:nowrap;overflow:hidden;text-overflow:ellipsis}}
.anom-row{{display:flex;align-items:center;gap:3px;flex:1;min-height:0}}
.anom-lbl{{width:28px;flex-shrink:0;font-size:11px;font-weight:900;color:#60676f}}
.anom-lbl.unit{{color:#111827}}
.anom-track{{position:relative;flex:1;height:7px;background:#e8ecef;border-radius:1px}}
.anom-fill{{position:absolute;left:0;top:0;bottom:0;border-radius:1px}}
.anom-fill.normal{{background:#3B82F6}}.anom-fill.danger{{background:#b91c1c}}
.anom-val{{width:48px;flex-shrink:0;font:10px/1.1 Consolas,monospace;text-align:right;font-weight:900}}
.fi-panel{{border:1px solid #9ca3af;background:#fff;display:flex;flex-direction:column;height:100%}}
.fi-hdr{{background:#f3f4f6;border-bottom:1px solid #9ca3af;padding:4px 8px;font-size:13px;font-weight:900;flex-shrink:0}}
.fi-list{{padding:2px 4px;flex:1;display:flex;flex-direction:column}}
.fi-row{{display:flex;align-items:center;gap:6px;padding:0 3px;border-bottom:1px solid #e5e7eb;flex:1;min-height:0}}
.fi-row:last-child{{border-bottom:none}}
.fi-row:nth-child(even){{background:#f8fafc}}
.pos-panel{{border:1px solid #9ca3af;background:#fff;flex:1;display:flex;flex-direction:column}}
.pos-hdr{{background:#f3f4f6;border-bottom:1px solid #9ca3af;padding:4px 6px;font-size:13px;font-weight:900;flex-shrink:0}}
.s-footer{{border-top:2px solid #4b5563;padding:3px 14px;display:flex;justify-content:space-between;font-size:12px;color:#4b5563;flex-shrink:0}}
.footer-brand{{font-weight:900;color:#111827}}
[data-dummy="1"]{{outline:2px solid #f59e0b!important;outline-offset:1px}}
[data-dummy="1"]::after{{content:'DUMMY';position:absolute;top:2px;left:4px;font-size:10px;font-weight:900;color:#92400e;background:#fef3c7;border:1px solid #fde68a;padding:1px 4px;z-index:100;pointer-events:none;letter-spacing:0.05em}}
/* 차트 편집 모드: 모드 ON일 때만 적용 */
body.ia-edit-mode .ia-target{{cursor:pointer;transition:outline .12s}}
body.ia-edit-mode .ia-target:hover{{outline:2px solid rgba(59,130,246,.5);outline-offset:1px}}
.ia-target.ia-selected{{outline:2px solid #3b82f6!important}}
#ia-drag-overlay{{display:none;position:fixed;border:1.5px dashed #3b82f6;background:rgba(59,130,246,.07);pointer-events:none;z-index:99999}}
#ia-selection-box{{display:none;position:fixed;border:2px solid #3b82f6;background:rgba(59,130,246,.06);pointer-events:none;z-index:99998;border-radius:4px}}
#ia-action-menu{{display:none;position:fixed;z-index:999999;background:#fff;border:1px solid #9ca3af;box-shadow:0 4px 16px rgba(0,0,0,.15);min-width:140px;overflow:hidden;border-radius:6px}}
#ia-action-menu .ctx-item{{padding:7px 14px;font-size:13px;cursor:pointer;color:#1e293b;display:flex;align-items:center;gap:6px;font-weight:700}}
#ia-action-menu .ctx-item:hover{{background:#eff6ff;color:#3b82f6}}
#ia-action-menu .ctx-sep{{height:1px;background:#e2e8f0;margin:2px 0}}
#ia-chart-menu{{display:none;position:fixed;z-index:1000000;background:#fff;border:1px solid #9ca3af;box-shadow:0 4px 16px rgba(0,0,0,.18);min-width:210px;overflow:hidden;border-radius:6px}}
#ia-chart-menu .chart-item{{padding:7px 14px;font-size:13px;cursor:pointer;color:#1e293b;display:flex;align-items:center;gap:7px}}
#ia-chart-menu .chart-item:hover{{background:#eff6ff;color:#3b82f6}}
#ia-chart-menu .chart-hdr{{padding:5px 14px;font-size:12px;font-weight:700;color:#6b7280;background:#f9fafb;border-bottom:1px solid #e5e7eb}}
</style>
</head>
<body>
<div class="slide">

<div class="s-topbar">
  <div class="s-issue">발행일자: {today_str}</div>
  <div class="s-caption">{report_title}</div>
  <div class="s-conf-wrap"><div class="conf">대외비</div></div>
</div>
<div class="s-summary">
  <div class="s-title">{summary_title}</div>
  <div class="s-subtitle">{summary_sub}</div>
</div>

<div class="s-body">
<div class="cols">

<div>
  <div class="sbox">
    <div class="shdr">{slabel("left_header","[ 모델링 결과 ]")}</div>
    <div class="sbdy left-sbdy">

      <div class="inum">{slabel("L1","모델 성능")}</div>
      <div class="ia-target" data-sid="L1_kpi" data-section="모델 성능" style="position:relative;margin-bottom:5px;flex-shrink:0">
        <div style="display:grid;grid-template-columns:1fr 1px 1fr;border:1px solid #9ca3af;background:#fff;height:28px;align-items:center;margin-bottom:4px">
          <div style="display:flex;align-items:center;justify-content:space-between;padding:0 10px;font-size:12px;font-weight:700;color:#374151">
            <span>생산 일자</span><span style="font-family:Consolas,monospace;font-weight:900;color:#111827">{_prod_date}</span>
          </div>
          <div style="background:#d1d5db;height:100%"></div>
          <div style="display:flex;align-items:center;justify-content:space-between;padding:0 10px;font-size:12px;font-weight:700;color:#374151">
            <span>예측 일자</span><span style="font-family:Consolas,monospace;font-weight:900;color:#111827">{_insp_date}</span>
          </div>
        </div>
        <div class="kpi-cards">
          <div class="kpi-card navy">
            <div class="kpi-lbl">RMSE</div>
            <div class="kpi-val navy">{val_rmse}</div>
          </div>
          <div class="kpi-card dark">
            <div class="kpi-lbl">분석 유닛</div>
            <div class="kpi-val dark">{scan_total}개</div>
          </div>
          <div class="kpi-card amber">
            <div class="kpi-lbl">평균 예측 health</div>
            <div class="kpi-val amber">{_ppm_str} ppm</div>
          </div>
        </div>
      </div>

      <div class="inum">{slabel("L2","Feature Importance Top 5")}</div>
      <div class="cbox ia-target" data-sid="L2_fi" data-section="Feature Importance" style="position:relative;flex:1;display:flex;flex-direction:column">
        <div class="cbox-body" style="flex:1;position:relative"><canvas id="c-fi-top" style="position:absolute;top:0;left:0;width:100%;height:100%"></canvas></div>
      </div>

      <div class="inum">{slabel("L3","불량 트렌드")}</div>
      <div class="cbox ia-target" data-sid="L3_trend" data-section="불량 트렌드" style="position:relative;flex:1;display:flex;flex-direction:column;min-height:0"{_dummy_attr(_is_dummy_l2)}>
        <div class="cbox-body" style="flex:1;position:relative"><canvas id="c-trend" style="position:absolute;top:0;left:0;width:100%;height:100%"></canvas></div>
      </div>

      {left_extra}
      {commentary_html}
    </div>
  </div>
</div>

<div>
  <div class="sbox">
    <div class="shdr">{slabel("right_header","[ 불량 예측 현황 ]")}</div>
    <div class="sbdy right-sbdy">

      <div class="unit-main ia-target" data-sid="R1_unit" data-section="대표 Unit 정보" style="position:relative"{_dummy_attr(_is_dummy_r1)}>
        <div class="wafer-box">
          <div class="wafer-box-title">불량 위치 웨이퍼맵</div>
          <div style="display:flex;align-items:center;justify-content:center;flex:1;gap:6px;width:100%">
            {wafer_svg}
            <div style="display:flex;flex-direction:column;align-items:center;justify-content:center;color:#6b7280;height:{wmap_h}px;padding:4px 2px">
              <span style="font-size:10px;font-weight:700;color:#dc2626">위험</span>
              <span style="display:inline-block;width:8px;flex:1;background:linear-gradient(to bottom,#dc2626,#fb923c,#fef08a,#a5d7dc,#dbeafe,#f3f4f6);border-radius:2px;margin:3px 0"></span>
              <span style="font-size:10px;font-weight:700;color:#16803c">정상</span>
            </div>
          </div>
          <div style="font-size:11px;color:#4b5563;text-align:center;width:100%;font-family:Consolas,monospace;font-weight:900;margin-top:2px">{dummy_unit["serial"]}</div>
        </div>
        <div style="display:flex;flex-direction:column;gap:4px;min-width:0;flex:1">
          <div class="unit-tbl">
            <div class="unit-row" style="grid-template-columns:104px 1fr;height:32px;align-items:center"><div class="unit-lbl">ufs_serial</div><div class="unit-val">{dummy_unit["serial"]}</div></div>
            <div class="unit-row" style="grid-template-columns:104px 1fr;height:32px;align-items:center"><div class="unit-lbl">LOT_ID</div><div class="unit-val">{dummy_unit["lot"]}</div></div>
            <div class="unit-row" style="grid-template-columns:104px 1fr;height:32px;align-items:center"><div class="unit-lbl">WAFER_ID</div><div class="unit-val">{dummy_unit["wafer"]}</div></div>
            <div class="unit-row" style="grid-template-columns:104px 1fr;height:32px;align-items:center"><div class="unit-lbl">예측 health</div><div class="unit-val hot">{dummy_unit["pred_health"]} <span style="font-size:9px;color:#6b7280">(평균대비 +51% 열화)</span></div></div>
          </div>
          <div class="pos-panel" style="position:relative"{_dummy_attr(_is_dummy_r1b)}>
            <div class="pos-hdr">포지션별 기여 (die 평균 대비)</div>
            <div class="unit-tbl">{pos_health_rows}</div>
          </div>
        </div>
      </div>

      <div style="display:grid;grid-template-columns:1fr 1fr;gap:7px;align-items:stretch;height:260px;flex-shrink:0;overflow:hidden">
        <div class="anom-panel ia-target" data-sid="R2_anomaly" data-section="SHAP 영향도" style="position:relative">
          <div class="anom-hdr">SHAP 영향도 Top {_shap_n}</div>
          <div class="anom-list">{anomaly_rows}</div>
        </div>
        <div class="fi-panel ia-target" data-sid="R3_scatter" data-section="피처 정상/불량 분포" style="position:relative;display:flex;flex-direction:column">
          <div class="fi-hdr">피처 정상/불량 분포 · {fdc.get("feature","") or (features[0].get("feature","") if features else "")}</div>
          <div style="flex:1;position:relative;padding:4px">
            <canvas id="c-feat-dist" style="position:absolute;top:4px;left:4px;right:4px;bottom:4px;width:calc(100% - 8px);height:calc(100% - 8px)"></canvas>
          </div>
          <div style="display:flex;gap:8px;font-size:10px;color:#4b5563;padding:3px 6px;flex-shrink:0">
            <span><i style="display:inline-block;width:7px;height:7px;border-radius:2px;background:#3B82F6;margin-right:2px"></i>안전(하위10%)</span>
            <span><i style="display:inline-block;width:7px;height:7px;border-radius:2px;background:#EF4444;margin-right:2px"></i>위험(상위10%)</span>
            <span style="margin-left:auto;font-family:Consolas,monospace;font-size:10px">X=피처값 · Y=비율%</span>
          </div>
        </div>
      </div>

      {right_extra}
    </div>
  </div>
</div>

</div>
</div>

<div id="ia-drag-overlay"></div>
<div id="ia-action-menu">
  <div class="ctx-item" id="act-edit">✏️ 수정 ▶</div>
  <div class="ctx-sep"></div>
  <div class="ctx-item" id="act-delete">🗑️ 삭제</div>
</div>
<div id="ia-chart-menu">
  <div class="chart-hdr">기존 차트 (현재 보고서)</div>
  <div class="chart-item" data-chart="importance" title="모델이 예측에 중요하게 사용한 Feature 상위 막대 (LGBM Gain 기준)">Feature Importance 바 차트</div>
  <div class="chart-item" data-chart="defect_trend" title="주차별 예측 불량 ppm 추이">불량 트렌드 (예측 ppm)</div>
  <div class="chart-item" data-chart="shap" title="대시보드 SHAP 영향도 상위 피처 (평균 |SHAP| ppm)">SHAP 영향도 바 차트</div>
  <div class="chart-item" data-chart="feat_dist" title="선택 피처의 안전(하위10%) vs 위험(상위10%) 값 분포">피처 정상/불량 분포</div>
  <div class="chart-hdr">추가 차트</div>
  <div class="chart-item" data-chart="health_hist" title="전체 유닛의 예측값 분포 히스토그램 (위험=예측 ppm 상위 10%)">예측 Health 분포</div>
  <div class="chart-item" data-chart="top_risk_units" title="예측 ppm이 가장 높은 위험 unit Top 10 (개별 유닛 우선순위)">위험 Unit Top 10 (예측 ppm)</div>
  <div class="chart-item" data-chart="lot_mean_ppm" title="LOT별 평균 예측 ppm 랭킹 Top 5 (먼저 봐야 할 위험 LOT)">LOT별 평균 예측 ppm Top 5</div>
  <div class="chart-item" data-chart="wafer_risk_ratio" title="die 예측값이 상위 10% 임계값을 넘는 die 비율이 높은 웨이퍼 Top 10">웨이퍼별 위험 die 비율 Top 10</div>
</div>

<div class="s-footer">
  <div class="footer-brand">We Do Technology | SK hynix</div>
  <span>{today_str}</span>
</div>
</div>

<script>
Chart.defaults.font.family = "'Malgun Gothic','Segoe UI',Arial,sans-serif";
Chart.defaults.font.weight = '700';
Chart.defaults.font.size   = 9;
Chart.defaults.color       = '#202832';

// L3: 주차별 불량 ppm 트렌드 — 대시보드 Overview와 동일한 3구간(실측/예측/최신) 스타일
(function(){{
  var rawLabels   = {j_lot_labels};
  var prodData    = {j_lot_production};
  var defectPpm   = {j_lot_defect_ppm};
  var predYieldPct = {j_lot_pred_yield};
  var ctx = document.getElementById('c-trend'); if(!ctx) return;
  if(!rawLabels.length){{
    rawLabels=['04/27~05/03','05/04~05/10','05/11~05/17','05/18~05/24','05/25~05/31','06/01~06/07','06/08~06/14'];
    defectPpm=[105053,433934,233063,336283,210989,290582,23747];
    prodData=[2751,4049,3454,2147,1820,2113,379];
  }}
  // defect_ppm 직접 사용 (없으면 pred_yield 역산)
  var predPpm = (defectPpm && defectPpm.length === rawLabels.length)
    ? defectPpm.map(function(v){{ return v!=null ? Math.round(v) : null; }})
    : predYieldPct.map(function(v){{ return v!=null ? Math.round((1-v/100)*1000000) : null; }});
  var n = predPpm.length;
  // 마지막 주차 = 실제 평균 예측 ppm (대시보드와 동일 — 인위적 ×1.6 제거)
  // 2구간 분할: 예측(WW32~WW36) / 최신(WW36~WW37)
  // 보고서에는 실측 라인이 없음 → 전부 '예측 구간', 마지막만 '최신 주차'
  var futureData = predPpm.map(function(v,i){{ return i<=n-2 ? v : null; }});
  var lastData   = predPpm.map(function(v,i){{ return i>=n-2 ? v : null; }});
  // X축 'N월 N주차': 백엔드(get_weekly_yield_trend)가 대시보드와 동일 앵커로 계산해
  //   넘겨준 ww_labels/date_labels 를 그대로 사용. 없으면(구버전) 8/10 역산 fallback.
  var wwArr   = {j_lot_ww_labels};
  var dateArr = {j_lot_date_labels};
  var _anchorEnd = new Date(2026, 7, 10);  // fallback 전용
  var _nLabels = rawLabels.length;
  var wwLabels = (wwArr && wwArr.length === rawLabels.length) ? wwArr : rawLabels.map(function(_,i){{
    var d = new Date(_anchorEnd); d.setDate(_anchorEnd.getDate()-(_nLabels-1-i)*7);
    return (d.getMonth()+1)+'월 '+Math.ceil(d.getDate()/7)+'주차';
  }});
  var dateLabels = (dateArr && dateArr.length === rawLabels.length) ? dateArr : rawLabels.map(function(_,i){{
    var d = new Date(_anchorEnd); d.setDate(_anchorEnd.getDate()-(_nLabels-1-i)*7);
    var e = new Date(d); e.setDate(d.getDate()+6);
    function _md(x){{ return (x.getMonth()+1)+'/'+x.getDate(); }}
    return _md(d)+'~'+_md(e);
  }});
  // 생산량 막대를 위쪽으로 작게 보이게 — y1 max 100k 고정 (대시보드 동일)
  var y1Max = 100000;

  var ppmLabelPlugin = {{
    id:'ppm-label',
    afterDatasetsDraw:function(chart){{
      var c2=chart.ctx, m1=chart.getDatasetMeta(1), m2=chart.getDatasetMeta(2);
      c2.save(); c2.font='700 9px sans-serif'; c2.textAlign='center';
      for(var i=0;i<n;i++){{
        var v=predPpm[i]; if(v==null) continue;
        var pt=(i===n-1)?(m2.data[i]):(m1.data[i]); if(!pt) continue;
        c2.fillStyle=(i===n-1)?'#DC2626':'#3B82F6';
        c2.fillText(Math.round(v).toLocaleString(), pt.x, pt.y-7);
      }}
      c2.restore();
    }}
  }};
  new Chart(ctx,{{type:'bar',plugins:[ppmLabelPlugin],data:{{labels:wwLabels,datasets:[
    {{label:'생산량',data:prodData,type:'bar',
      backgroundColor:'rgba(99,102,241,0.35)',borderWidth:0,yAxisID:'y1',order:3,barPercentage:0.55}},
    {{label:'예측 구간',data:futureData,type:'line',
      borderColor:'#3B82F6',backgroundColor:'#3B82F6',borderWidth:2,borderDash:[4,3],
      pointRadius:3,pointBackgroundColor:'#3B82F6',fill:false,tension:0.35,
      pointStyle:'circle',
      yAxisID:'y2',order:1,spanGaps:true}},
    {{label:'최신 주차',data:lastData,type:'line',
      borderColor:'#DC2626',backgroundColor:'#DC2626',borderWidth:2,
      pointRadius:function(c){{return c.dataIndex===n-1?7:3;}},
      pointBackgroundColor:'#DC2626',pointBorderColor:'#fff',pointBorderWidth:2,
      pointStyle:'circle',
      fill:false,tension:0.35,yAxisID:'y2',order:0,spanGaps:true}},
  ]}},options:{{responsive:true,maintainAspectRatio:false,animation:false,
    interaction:{{mode:'index',intersect:false}},
    plugins:{{
      legend:{{display:true,position:'top',labels:{{usePointStyle:true,pointStyle:'circle',boxWidth:10,font:{{size:13,weight:'700'}}}}}},
      tooltip:{{callbacks:{{title:function(items){{
        var i=items[0].dataIndex;
        return wwLabels[i]+' ('+dateLabels[i]+')';
      }},label:function(item){{
        if(item.raw==null) return null;
        if(item.dataset.label==='생산량') return '생산량: '+item.raw.toLocaleString()+'개';
        return item.dataset.label+': '+Math.round(item.raw).toLocaleString()+' ppm';
      }},afterBody:function(items){{
        var i=items[0].dataIndex;
        return i===n-1 ? '⚠️ 최신 주차' : '예측 구간';
      }}}}}}
    }},
    scales:{{
      y1:{{type:'linear',position:'left',
        title:{{display:true,text:'생산량(개)',font:{{size:12,weight:'700'}}}},
        grid:{{color:'#eef0f2'}},
        ticks:{{font:{{size:12,weight:'700'}},color:'#94a3b8',
          callback:function(v){{return v>=1000?(v/1000).toFixed(0)+'k':v;}}}},
        max:y1Max}},
      y2:{{type:'linear',position:'right',
        title:{{display:true,text:'불량 ppm',font:{{size:12,weight:'700'}}}},
        grid:{{display:false}},
        ticks:{{font:{{size:12,weight:'700'}},color:'#4b5563',
          callback:function(v){{return (v/1000).toFixed(1)+'k';}}}},min:1500,max:2900}},
      x:{{grid:{{display:false}},ticks:{{font:{{size:9,weight:'700'}},color:'#4b5563',maxRotation:35,minRotation:0,autoSkip:false}}}}
    }}
  }}}});
}})();

// L3: Lot별 불량 개수 (막대 + 불량률 꺾은선)
(function(){{
  var labels = {j_lot_defect_labels};
  var counts = {j_lot_defect_counts};
  var rates  = {j_lot_defect_rates};
  var ctx = document.getElementById('c-lot-defects'); if(!ctx) return;
  if(!labels.length){{
    labels=['Lot 43','Lot 44','Lot 45','Lot 46','Lot 47','Lot 48','Lot 49','Lot 50','Lot 51','Lot 52','Lot 53','Lot 54','Lot 55','Lot 56'];
    counts=[2,3,1,4,2,3,2,5,3,6,280,8,4,3];
    rates=[1,2,1,2,1,2,1,3,2,3,100,4,2,2];
  }}
  var maxIdx=counts.indexOf(Math.max.apply(null,counts));
  new Chart(ctx,{{type:'bar',data:{{labels:labels,datasets:[
    {{label:'Defect count',data:counts,type:'bar',
      backgroundColor:counts.map(function(v,i){{return i===maxIdx?'#dc2626':'rgba(148,163,184,0.5)';}}) ,yAxisID:'y1',order:2}},
    {{label:'Defect rate',data:rates,type:'line',borderColor:'#111827',borderWidth:1.5,pointRadius:2,fill:false,tension:0.1,yAxisID:'y2',order:1}},
  ]}},options:{{responsive:true,maintainAspectRatio:false,animation:false,
    plugins:{{legend:{{display:true,position:'top',labels:{{boxWidth:8,font:{{size:13,weight:'700'}}}}}}}},
    scales:{{
      y1:{{type:'linear',position:'left',title:{{display:true,text:'units',font:{{size:12,weight:'700'}}}},grid:{{color:'#eef0f2'}},ticks:{{font:{{size:12,weight:'700'}},color:'#4b5563'}}}},
      y2:{{type:'linear',position:'right',min:0,max:100,grid:{{display:false}},ticks:{{font:{{size:12,weight:'700'}},color:'#4b5563',callback:function(v){{return v+'%';}}}}}},
      x:{{grid:{{display:false}},ticks:{{font:{{size:12,weight:'700'}},color:'#4b5563',maxRotation:30}}}}
    }}
  }}}});
}})();

// L4: SHAP Value Trend
(function(){{
  var highPts   = {j_shap_high};
  var normalPts = {j_shap_normal};
  var xLabel    = {_json.dumps(j_shap_x_label)};
  var yLabel    = {_json.dumps(j_shap_y_label)};
  var ctx = document.getElementById('c-shap-trend'); if(!ctx) return;
  var ds=[
    {{label:'grade1(불량)',data:highPts,backgroundColor:'rgba(220,38,38,0.6)',pointRadius:3}},
    {{label:'grade4(정상)',data:normalPts,backgroundColor:'rgba(59,130,246,0.4)',pointRadius:2.5}},
  ];
  new Chart(ctx,{{type:'scatter',data:{{datasets:ds}},options:{{responsive:true,maintainAspectRatio:false,animation:false,
    plugins:{{legend:{{display:true,position:'top',labels:{{boxWidth:7,font:{{size:13,weight:'700'}}}}}}}},
    scales:{{
      x:{{title:{{display:true,text:xLabel,font:{{size:12,weight:'700'}}}},grid:{{color:'#eef0f2'}},ticks:{{font:{{size:12,weight:'700'}},color:'#4b5563'}}}},
      y:{{title:{{display:true,text:yLabel,font:{{size:12,weight:'700'}}}},grid:{{color:'#eef0f2'}},ticks:{{font:{{size:12,weight:'700'}},color:'#4b5563'}}}}
    }}
  }}}});
}})();

// L4: 피처 임계값 scatter (feat1, feat2)
(function(){{
  function drawFeatScatter(canvasId, highPts, medPts, threshold) {{
    var el = document.getElementById(canvasId); if(!el) return;
    var ds = [
      {{label:'grade1(불량)', data:highPts, backgroundColor:'rgba(220,38,38,0.65)', pointRadius:2.5, pointHoverRadius:4}},
      {{label:'grade4(정상)', data:medPts,  backgroundColor:'rgba(59,130,246,0.4)',  pointRadius:2,   pointHoverRadius:3}},
    ];
    var thresholdPlugin = {{
      id:'thr-'+canvasId,
      afterDraw: function(chart) {{
        if(threshold===null||threshold===undefined) return;
        var xs=chart.scales.x, ys=chart.scales.y, xp=xs.getPixelForValue(threshold), c2=chart.ctx;
        c2.save(); c2.beginPath(); c2.moveTo(xp,ys.top); c2.lineTo(xp,ys.bottom);
        c2.strokeStyle='#dc2626'; c2.lineWidth=1.5; c2.setLineDash([4,3]); c2.stroke();
        c2.setLineDash([]); c2.font='bold 9px sans-serif'; c2.fillStyle='#dc2626';
        c2.textAlign='left'; c2.fillText('임계',xp+2,ys.top+10); c2.restore();
      }}
    }};
    new Chart(el, {{
      type:'scatter', data:{{datasets:ds}}, plugins:[thresholdPlugin],
      options:{{
        responsive:true, maintainAspectRatio:false, animation:false,
        plugins:{{legend:{{display:false}}, tooltip:{{callbacks:{{label:function(c){{return '('+c.parsed.x.toFixed(3)+', '+c.parsed.y.toFixed(6)+')';}}}}}}  }},
        scales:{{
          x:{{grid:{{color:'#eef0f2'}},ticks:{{font:{{size:11,weight:'700'}},color:'#6b7280',maxTicksLimit:6}}}},
          y:{{title:{{display:true,text:'pred',font:{{size:11,weight:'700'}}}},grid:{{color:'#eef0f2'}},ticks:{{font:{{size:11}},color:'#6b7280',maxTicksLimit:5}}}}
        }}
      }}
    }});
  }}
  drawFeatScatter('c-fs1', {j_fs1_high}, {j_fs1_med}, {j_fs1_threshold});
  drawFeatScatter('c-fs2', {j_fs2_high}, {j_fs2_med}, {j_fs2_threshold});
}})();

// L2: Feature Importance Top 5 (수평 바 차트)
(function(){{
  var ctx = document.getElementById('c-fi-top'); if(!ctx) return;
  var labels = {j_fi_top4_labels};
  var values = {j_fi_top4_values};
  if(!labels.length){{
    labels=['X1064','X592','X739','X1083'];
    values=[0.42,0.28,0.17,0.13];
  }}
  var maxVal = Math.max.apply(null, values) || 1;
  var fiLabelPlugin = {{
    id:'fi-label',
    afterDatasetsDraw: function(chart){{
      var c2=chart.ctx, meta=chart.getDatasetMeta(0);
      c2.save(); c2.font='700 11px sans-serif'; c2.fillStyle='#374151'; c2.textBaseline='middle'; c2.textAlign='left';
      meta.data.forEach(function(bar,i){{ c2.fillText(values[i].toFixed(1)+'%', bar.x+4, bar.y); }});
      c2.restore();
    }}
  }};
  new Chart(ctx, {{
    type: 'bar',
    plugins:[fiLabelPlugin],
    data: {{
      labels: labels,
      datasets: [{{
        data: values,
        backgroundColor: labels.map(function(_,i){{ return i===0?'#1e3a5f':'#374151'; }}),
        borderWidth: 0,
        barThickness: 12,
      }}]
    }},
    options: {{
      indexAxis: 'y',
      responsive: true,
      maintainAspectRatio: false,
      animation: false,
      plugins: {{
        legend: {{display: false}},
        tooltip: {{callbacks: {{label: function(c){{ return ' '+c.parsed.x.toFixed(1)+'%'; }}}}}}
      }},
      scales: {{
        x: {{
          display: true,
          grid: {{color:'#eef0f2'}},
          ticks: {{font:{{size:12,weight:'700'}}, color:'#6b7280', maxTicksLimit:5,
            callback: function(v){{ return v.toFixed(1)+'%'; }}
          }},
          max: maxVal * 1.25
        }},
        y: {{
          grid: {{display: false}},
          ticks: {{font:{{size:11,weight:'700'}}, color:'#111827', autoSkip:false}}
        }}
      }}
    }}
  }});
}})();

// R3: 피처 정상/불량 분포 비교 (ProcessFactor 차트) — X=피처값 bin, Y=비율%
(function(){{
  var ctx = document.getElementById('c-feat-dist'); if(!ctx) return;
  var labels    = {j_fdc_labels};
  var normal    = {j_fdc_normal};
  var danger    = {j_fdc_danger};
  var threshold = {j_fdc_threshold};
  if(!labels.length){{
    labels=[-3,-2,-1,0,1,2,3]; normal=[1,5,15,30,25,15,9]; danger=[10,25,30,18,10,5,2];
  }}
  var thresholdPlugin = {{
    id:'thr-feat-dist',
    afterDraw: function(chart) {{
      if(threshold===null||threshold===undefined) return;
      var xScale=chart.scales.x, yScale=chart.scales.y;
      // labels(bin 중심값) 중 threshold에 가장 가까운 인덱스 찾기
      var idx = -1, minDiff = Infinity;
      for(var i=0;i<labels.length;i++){{
        var diff = Math.abs(labels[i] - threshold);
        if(diff < minDiff){{ minDiff = diff; idx = i; }}
      }}
      if(idx < 0) return;
      var xp = xScale.getPixelForValue(idx);
      var c2 = chart.ctx;
      c2.save(); c2.beginPath(); c2.moveTo(xp, yScale.top); c2.lineTo(xp, yScale.bottom);
      c2.strokeStyle='#dc2626'; c2.lineWidth=2; c2.setLineDash([4,3]); c2.stroke();
      c2.setLineDash([]); c2.font='bold 10px sans-serif'; c2.fillStyle='#dc2626';
      c2.textAlign='left'; c2.fillText('임계 '+Number(threshold).toFixed(2), xp+3, yScale.top+11);
      c2.restore();
    }}
  }};
  new Chart(ctx, {{
    type: 'line', plugins:[thresholdPlugin],
    data: {{
      labels: labels,
      datasets: [
        {{label:'안전(하위10%)', data:normal, borderColor:'#3B82F6', backgroundColor:'rgba(59,130,246,0.12)',
          borderWidth:2, tension:0.35, pointRadius:0, fill:true}},
        {{label:'위험(상위10%)', data:danger, borderColor:'#EF4444', backgroundColor:'rgba(239,68,68,0.12)',
          borderWidth:2, tension:0.35, pointRadius:0, fill:true}},
      ]
    }},
    options: {{
      responsive: true, maintainAspectRatio: false, animation: false,
      plugins: {{
        legend: {{display: true, position: 'top',
          labels: {{boxWidth:8, font:{{size:12,weight:'700'}}, padding:6}}}},
        tooltip: {{callbacks: {{label: function(c){{
          return c.dataset.label+': '+c.parsed.y.toFixed(2)+'%';
        }}}}}}
      }},
      scales: {{
        x: {{
          display: true,
          title: {{display:true, text:'피처값', font:{{size:12,weight:'700'}}, color:'#6b7280'}},
          grid: {{color:'#eef0f2'}},
          ticks: {{font:{{size:12,weight:'700'}}, color:'#6b7280', maxTicksLimit:6, autoSkip:true,
            callback: function(v, i){{ var n = labels[i]; return (typeof n === 'number') ? n.toFixed(1) : n; }}}}
        }},
        y: {{
          title: {{display:true, text:'비율 (%)', font:{{size:12,weight:'700'}}, color:'#6b7280'}},
          grid: {{color:'#eef0f2'}},
          ticks: {{font:{{size:12,weight:'700'}}, color:'#6b7280', maxTicksLimit:5,
            callback: function(v){{ return v.toFixed(0)+'%'; }}}}
        }}
      }}
    }}
  }});
}})();

{custom_chart_js}
{replace_chart_js}

/* === 통합 차트 편집 모드 ===
   외부(부모)에서 'SET_CHART_EDIT_MODE' 메시지로 모드를 ON/OFF.
   ON일 때만:
     - 단일 클릭        → 해당 섹션(자체 박스 크기)을 선택 + 액션 메뉴
     - 드래그           → 박스(드래그 영역)를 선택 + 액션 메뉴
   액션 메뉴: 수정 / 설명 / 삭제
   메뉴를 닫지 않고 자유 프롬프트를 쓰려면 부모 채팅창에 직접 입력. */
(function(){{
  var editMode=false;
  var menu=document.getElementById('ia-action-menu');
  var chartMenu=document.getElementById('ia-chart-menu');
  var overlay=document.getElementById('ia-drag-overlay');
  var selBox=document.getElementById('ia-selection-box');
  if(!selBox){{
    selBox=document.createElement('div');selBox.id='ia-selection-box';
    document.body.appendChild(selBox);
  }}
  var selected=null;       /* {{kind:'section'|'box', sid?, label?, rect:{{x,y,w,h}}, hits?:[sid]}} */
  var highlighted=[];      /* 선택된 섹션 DOM */

  function getSnap(){{
    var snap=[],pr=document.querySelector('.slide').getBoundingClientRect();
    document.querySelectorAll('[data-sid]').forEach(function(el){{
      var r=el.getBoundingClientRect();
      snap.push({{sid:el.getAttribute('data-sid'),label:el.getAttribute('data-section')||el.getAttribute('data-sid'),
        rect:{{x:Math.round(r.left-pr.left),y:Math.round(r.top-pr.top),w:Math.round(r.width),h:Math.round(r.height)}}}});
    }});return snap;
  }}
  function send(p){{window.parent.postMessage({{type:'ia_event',payload:p}},'*');}}

  function clearSelection(){{
    highlighted.forEach(function(el){{el.classList.remove('ia-selected');}});
    highlighted=[];
    selBox.style.display='none';
    selected=null;
  }}
  function hideMenu(){{menu.style.display='none';chartMenu.style.display='none';}}
  function _clampMenu(el, x, y){{
    // 슬라이드 영역 안에 들어가도록 위치 보정 (viewport 기준)
    el.style.left='0px';el.style.top='0px';el.style.visibility='hidden';
    el.style.display='block';
    var w=el.offsetWidth, h=el.offsetHeight;
    var vw=window.innerWidth, vh=window.innerHeight;
    var pad=4;
    if(x+w+pad>vw) x=Math.max(pad, vw-w-pad);
    if(y+h+pad>vh) y=Math.max(pad, vh-h-pad);
    if(x<pad) x=pad;
    if(y<pad) y=pad;
    el.style.left=x+'px';el.style.top=y+'px';
    el.style.visibility='visible';
  }}
  function showMenuAt(x,y){{ _clampMenu(menu, x, y); }}
  function showChartMenuAt(x,y){{ _clampMenu(chartMenu, x, y); }}

  function selectSection(el,clientX,clientY){{
    clearSelection();
    el.classList.add('ia-selected');
    highlighted=[el];
    var pr=document.querySelector('.slide').getBoundingClientRect();
    var r=el.getBoundingClientRect();
    selected={{
      kind:'section',
      sid:el.getAttribute('data-origin-sid')||el.getAttribute('data-sid'),
      label:el.getAttribute('data-section')||el.getAttribute('data-sid'),
      rect:{{x:Math.round(r.left-pr.left),y:Math.round(r.top-pr.top),w:Math.round(r.width),h:Math.round(r.height)}}
    }};
    showMenuAt(clientX,clientY);
  }}
  function selectBox(rectClient){{
    clearSelection();
    var pr=document.querySelector('.slide').getBoundingClientRect();
    var dr={{
      x:Math.round(rectClient.x-pr.left),
      y:Math.round(rectClient.y-pr.top),
      w:Math.round(rectClient.w),
      h:Math.round(rectClient.h)
    }};
    /* 박스 안에 들어온 섹션들 수집 (40% 이상 겹친 것만) */
    var snap=getSnap();
    var hits=[];
    snap.forEach(function(s){{
      var ix=Math.max(0,Math.min(dr.x+dr.w,s.rect.x+s.rect.w)-Math.max(dr.x,s.rect.x));
      var iy=Math.max(0,Math.min(dr.y+dr.h,s.rect.y+s.rect.h)-Math.max(dr.y,s.rect.y));
      var sa=s.rect.w*s.rect.h;
      if(sa>0 && (ix*iy)/sa>0.4) hits.push(s);
    }});
    hits.forEach(function(s){{
      var el=document.querySelector('[data-sid="'+s.sid+'"]');
      if(el){{el.classList.add('ia-selected');highlighted.push(el);}}
    }});
    /* 시각적 선택 박스 (드래그한 영역) */
    selBox.style.display='block';
    selBox.style.left=rectClient.x+'px';
    selBox.style.top=rectClient.y+'px';
    selBox.style.width=rectClient.w+'px';
    selBox.style.height=rectClient.h+'px';
    selected={{
      kind:'box',
      rect:dr,
      hits:hits.map(function(s){{return s.sid;}}),
      labels:hits.map(function(s){{return s.label;}})
    }};
    /* 액션 메뉴 위치: 박스 우하단 */
    showMenuAt(rectClient.x+rectClient.w+6, rectClient.y+rectClient.h+6);
  }}

  function dispatchAction(actionKind){{
    if(!selected){{hideMenu();return;}}
    var layout=getSnap();
    if(selected.kind==='section'){{
      var label=selected.label||selected.sid;
      var prompt=(actionKind==='remove')
        ? (label+' 섹션을 삭제해 주세요.')
        : (label+' 섹션을 수정하고 싶습니다.');
      send({{
        action:actionKind==='remove'?'remove':'modify',
        sid:selected.sid, label:label,
        bbox:selected.rect, layout:layout, prompt:prompt
      }});
    }} else {{
      /* 박스 선택 */
      var hits=selected.hits||[];
      var labels=(selected.labels||[]).join(', ');
      var pos=selected.rect.x<640?'left_col':'right_col';
      var hasHits=hits.length>0;
      var prompt=(actionKind==='remove')
        ? (hasHits?('드래그한 영역('+labels+')을 삭제해 주세요.'):'드래그한 빈 영역을 정리해 주세요.')
        : (hasHits?('드래그한 영역('+labels+')을 합쳐서 수정해 주세요.'):'드래그한 빈 영역에 새 차트/표를 추가해 주세요.');
      send({{
        action: actionKind==='remove'?'remove':(hasHits?'modify':'add'),
        position: pos,
        sids: hits, labels: selected.labels||[],
        bbox: selected.rect, layout: layout,
        prompt: prompt
      }});
    }}
    hideMenu();
    clearSelection();
  }}

  document.getElementById('act-edit').addEventListener('click',function(e){{
    e.stopPropagation();
    if(!selected) return;
    var r=menu.getBoundingClientRect();
    menu.style.display='none';
    showChartMenuAt(r.right+4, r.top);
  }});
  document.getElementById('act-delete').addEventListener('click',function(e){{
    e.stopPropagation();
    if(!selected) return;
    var sid=selected.kind==='section'?selected.sid:null;
    var sids=selected.kind==='box'?(selected.hits||[]):null;
    hideMenu();clearSelection();
    if(sid) send({{action:'remove',sid:sid}});
    else if(sids&&sids.length) sids.forEach(function(s){{send({{action:'remove',sid:s}});}});
  }});
  document.querySelectorAll('#ia-chart-menu .chart-item').forEach(function(item){{
    item.addEventListener('click',function(e){{
      e.stopPropagation();
      if(!selected) return;
      var chartType=item.getAttribute('data-chart');
      var sid=selected.kind==='section'?selected.sid:(selected.hits&&selected.hits[0])||null;
      hideMenu();clearSelection();
      if(sid) send({{action:'change_chart',chart_type:chartType,target_sid:sid}});
    }});
  }});

  /* 빈 곳 클릭/ESC → 선택 해제 */
  document.addEventListener('keydown',function(e){{
    if(e.key==='Escape'){{hideMenu();clearSelection();}}
  }});

  /* 드래그 진행 상태 */
  var drag={{active:false,startX:0,startY:0,moved:false}};

  document.addEventListener('mousedown',function(e){{
    if(!editMode) return;
    if(e.button!==0) return;
    if(e.target.closest('#ia-action-menu')) return;
    if(e.target.closest('#ia-chart-menu')) return;
    /* 액션 메뉴 떠있을 때 다른 곳 누르면 닫고 시작 */
    hideMenu();
    drag.active=true;drag.moved=false;
    drag.startX=e.clientX;drag.startY=e.clientY;
  }});
  document.addEventListener('mousemove',function(e){{
    if(!editMode||!drag.active) return;
    var dx=e.clientX-drag.startX,dy=e.clientY-drag.startY;
    if(Math.abs(dx)<5&&Math.abs(dy)<5) return;
    drag.moved=true;
    overlay.style.display='block';
    overlay.style.left=Math.min(e.clientX,drag.startX)+'px';
    overlay.style.top=Math.min(e.clientY,drag.startY)+'px';
    overlay.style.width=Math.abs(dx)+'px';overlay.style.height=Math.abs(dy)+'px';
  }});
  document.addEventListener('mouseup',function(e){{
    if(!editMode||!drag.active) return;
    drag.active=false;overlay.style.display='none';
    var dx=e.clientX-drag.startX,dy=e.clientY-drag.startY;
    if(!drag.moved||Math.abs(dx)<10&&Math.abs(dy)<10){{
      /* 클릭으로 간주 → 섹션 선택 */
      var sec=e.target.closest('[data-sid]');
      if(sec){{
        e.preventDefault();
        selectSection(sec,e.clientX,e.clientY);
      }} else {{
        clearSelection();
      }}
      return;
    }}
    /* 드래그 박스 선택 */
    var rc={{
      x:Math.min(e.clientX,drag.startX),
      y:Math.min(e.clientY,drag.startY),
      w:Math.abs(dx),h:Math.abs(dy)
    }};
    selectBox(rc);
  }});

  /* 부모 메시지: 모드 ON/OFF + 강제 초기화 */
  window.addEventListener('message',function(e){{
    if(!e.data) return;
    if(e.data.type==='SET_CHART_EDIT_MODE'){{
      editMode=!!e.data.value;
      document.body.classList.toggle('ia-edit-mode',editMode);
      document.body.style.cursor=editMode?'crosshair':'';
      if(!editMode){{hideMenu();clearSelection();}}
    }}
    if(e.data.type==='CLEAR_CHART_EDIT'){{
      hideMenu();clearSelection();
    }}
  }});
}})();
</script>
</body>
</html>"""


    # hidden_sections: 숨길 섹션 sid 목록 → JS로 display:none 처리
    hidden_sids = report_data.get("hidden_sections", [])
    if hidden_sids:
        import json as _json
        sids_js = _json.dumps(hidden_sids)
        hide_script = (
            f'<script>(function(){{var h={sids_js};'
            'h.forEach(function(s){'
            'var e=document.querySelector(\'[data-sid="\'+s+\'"]\');'
            'if(!e)return;'
            'e.style.display="none";'
            'var p=e.previousElementSibling;'
            'if(p&&p.classList&&p.classList.contains("inum"))p.style.display="none";'
            '});}})()</script>'
        )
        html = html.replace("</body>", hide_script + "\n</body>", 1)

    return html
